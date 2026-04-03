"""
vault_ai_gui.py
Vault AI + TheHelper browser (Chromium via Qt WebEngine)
+ Document Upload (PDF, PPTX, DOCX)

REQUIREMENTS:
  pip install PySide6 requests
  pip install PyPDF2 python-pptx python-docx
  pip install faiss-cpu sentence-transformers ollama
"""

import os
import re
import sys
import json
import pickle
import shutil
import threading
import subprocess
import webbrowser
import importlib.util
import tkinter as tk
from tkinter import (
    ttk, scrolledtext, messagebox, filedialog)
from datetime import datetime
import urllib.parse
import urllib.request
import urllib.error

# -- Optional heavy deps --
try:
    import faiss
    import numpy as np
    from sentence_transformers import SentenceTransformer
    import ollama
    VAULT_AVAILABLE = True
except ImportError:
    VAULT_AVAILABLE = False

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

# -- Document parsing deps --
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

HAS_QT_BROWSER = (
    importlib.util.find_spec("PySide6") is not None
    or importlib.util.find_spec("PyQt6") is not None
)

# ==================================================
#  CONFIG
# ==================================================
NOTES_FOLDER = os.path.expanduser(
    r"C:\Users\ASUS\Obsidian")
UPLOADS_DIR  = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "uploaded_docs")
INDEX_FILE   = "faiss_index.bin"
META_FILE    = "metadata.pkl"
DOCS_LOG     = "uploaded_docs.json"
LOG_FILE     = "chat_logs.json"
CHUNK_SIZE   = 800
TOP_K        = 5
MODEL        = "llama3"
EMBED_MODEL  = "all-MiniLM-L6-v2"
MAX_TOKENS   = 2000
BASE_URL     = "https://thehelpers.vercel.app"
DOWNLOAD_DIR = os.path.join(
    os.path.expanduser("~"),
    "Downloads", "TheHelper")

# ==================================================
#  COLOUR PALETTE
# ==================================================
BG        = "#0f0f0f"
BG2       = "#1a1a1a"
BG3       = "#222222"
ACCENT    = "#e8c547"
ACCENT2   = "#4fc3f7"
FG        = "#e0e0e0"
FG_DIM    = "#777777"
RED       = "#ef5350"
GREEN     = "#66bb6a"
BORDER    = "#333333"
FONT_MONO = ("Consolas", 10)
FONT_UI   = ("Segoe UI", 10)
FONT_H1   = ("Segoe UI", 14, "bold")

# ==================================================
#  DOCUMENT PARSING
# ==================================================

def extract_pdf(filepath):
    if not HAS_PDF:
        return ""
    text = ""
    try:
        reader = PyPDF2.PdfReader(filepath)
        for page in reader.pages:
            pt = page.extract_text()
            if pt:
                text += pt + "\n"
    except Exception as exc:
        text = "[PDF error: " + str(exc) + "]"
    return text


def extract_pptx(filepath):
    if not HAS_PPTX:
        return ""
    text = ""
    try:
        prs = Presentation(filepath)
        for snum, slide in enumerate(
                prs.slides, 1):
            text += "\n--- Slide " + str(
                snum) + " ---\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in (
                            shape.text_frame
                            .paragraphs):
                        line = para.text.strip()
                        if line:
                            text += line + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = []
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct:
                                cells.append(ct)
                        if cells:
                            text += (
                                " | ".join(cells)
                                + "\n")
            if slide.has_notes_slide:
                notes = (
                    slide.notes_slide
                    .notes_text_frame
                    .text.strip())
                if notes:
                    text += (
                        "Notes: " + notes + "\n")
    except Exception as exc:
        text = "[PPTX error: " + str(exc) + "]"
    return text


def extract_docx(filepath):
    if not HAS_DOCX:
        return ""
    text = ""
    try:
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            line = para.text.strip()
            if line:
                text += line + "\n"
        for table in doc.tables:
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    ct = cell.text.strip()
                    if ct:
                        cells.append(ct)
                if cells:
                    text += (
                        " | ".join(cells) + "\n")
    except Exception as exc:
        text = "[DOCX error: " + str(exc) + "]"
    return text


def extract_text(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext == ".pptx":
        return extract_pptx(filepath)
    elif ext in (".docx", ".doc"):
        return extract_docx(filepath)
    elif ext in (".md", ".txt"):
        try:
            return open(
                filepath, encoding="utf-8",
                errors="ignore").read()
        except Exception:
            return ""
    return ""


def get_supported_types():
    types = [("Markdown", "*.md"),
             ("Text", "*.txt")]
    if HAS_PDF:
        types.append(("PDF", "*.pdf"))
    if HAS_PPTX:
        types.append(("PowerPoint", "*.pptx"))
    if HAS_DOCX:
        types.append(("Word", "*.docx"))
    all_exts = " ".join(t[1] for t in types)
    return [("All Supported", all_exts)] + types


def get_dep_status():
    lines = []
    lines.append(
        "PDF  (PyPDF2):       "
        + ("installed" if HAS_PDF
           else "MISSING - pip install PyPDF2"))
    lines.append(
        "PPTX (python-pptx):  "
        + ("installed" if HAS_PPTX
           else "MISSING - pip install python-pptx"))
    lines.append(
        "DOCX (python-docx):  "
        + ("installed" if HAS_DOCX
           else "MISSING - pip install python-docx"))
    lines.append(
        "FAISS + Ollama:      "
        + ("installed" if VAULT_AVAILABLE
           else "MISSING"))
    return "\n".join(lines)


# ==================================================
#  UPLOADED DOCS TRACKING
# ==================================================

def load_docs_log():
    if os.path.exists(DOCS_LOG):
        try:
            return json.load(open(
                DOCS_LOG, "r", encoding="utf-8"))
        except Exception:
            pass
    return []


def save_docs_log(docs):
    json.dump(docs, open(
        DOCS_LOG, "w", encoding="utf-8"), indent=2)


# ==================================================
#  GOOGLE DRIVE HELPERS
# ==================================================

def gdrive_id_from_url(url):
    for p in [r"/file/d/([a-zA-Z0-9_-]+)",
              r"id=([a-zA-Z0-9_-]+)",
              r"/d/([a-zA-Z0-9_-]+)/"]:
        m = re.search(p, url)
        if m:
            return m.group(1)
    return None


def make_direct_download(url):
    fid = gdrive_id_from_url(url)
    if fid:
        return ("https://drive.google.com/"
                "uc?export=download&id=" + fid)
    return url


def download_file(url, dest_dir, label,
                  log_fn=None):
    os.makedirs(dest_dir, exist_ok=True)
    if "drive.google.com" in url:
        dl_url = make_direct_download(url)
    else:
        dl_url = url
    safe = re.sub(r'[\\/:*?"<>|]', "_", label)
    if not safe.endswith(".pdf"):
        safe += ".pdf"
    dest = os.path.join(dest_dir, safe)
    if log_fn:
        log_fn("Downloading: " + safe)
    try:
        if REQUESTS_AVAILABLE:
            r = requests.get(
                dl_url, stream=True, timeout=30,
                headers={
                    "User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            with open(dest, "wb") as f:
                for chunk in r.iter_content(8192):
                    f.write(chunk)
        else:
            req = urllib.request.Request(
                dl_url,
                headers={
                    "User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(
                    req, timeout=30) as r:
                with open(dest, "wb") as f:
                    f.write(r.read())
        if log_fn:
            log_fn("Saved: " + dest)
        return dest
    except Exception as exc:
        if log_fn:
            log_fn("Failed: " + str(exc))
        return None


# ==================================================
#  VAULT AI CORE
# ==================================================
embed_model = None
index       = None
metadata    = None


def load_vault(log_fn=None):
    global embed_model, index, metadata
    if not VAULT_AVAILABLE:
        if log_fn:
            log_fn("Vault deps missing.")
        return False
    embed_model = SentenceTransformer(EMBED_MODEL)
    index, metadata = _build_or_load(log_fn)
    return index is not None


def _build_or_load(log_fn=None):
    if (os.path.exists(INDEX_FILE)
            and os.path.exists(META_FILE)):
        if log_fn:
            log_fn("Loading existing FAISS index...")
        idx = faiss.read_index(INDEX_FILE)
        meta = pickle.load(open(META_FILE, "rb"))
        if log_fn:
            log_fn("Loaded " + str(idx.ntotal)
                   + " vectors.")
        return idx, meta
    return rebuild_full_index(log_fn)


def _chunk_text(text, filename, filepath, source):
    chunks = []
    for i in range(0, len(text), CHUNK_SIZE):
        chunk = text[i:i + CHUNK_SIZE].strip()
        if chunk:
            meta = {
                "text": chunk,
                "path": filepath,
                "file": filename,
                "source": source,
                "chunk_idx": i // CHUNK_SIZE
            }
            chunks.append((chunk, meta))
    return chunks


def rebuild_full_index(log_fn=None):
    global index, metadata
    if log_fn:
        log_fn("Building full index...")
    docs, meta = [], []

    # 1. Obsidian notes
    if os.path.exists(NOTES_FOLDER):
        if log_fn:
            log_fn("Scanning: " + NOTES_FOLDER)
        for root_dir, _, files in os.walk(
                NOTES_FOLDER):
            for f in files:
                if f.endswith(".md"):
                    path = os.path.join(
                        root_dir, f)
                    try:
                        text = open(
                            path,
                            encoding="utf-8",
                            errors="ignore"
                        ).read()
                        chunks = _chunk_text(
                            text, f, path,
                            "obsidian")
                        docs.extend(
                            c[0] for c in chunks)
                        meta.extend(
                            c[1] for c in chunks)
                    except Exception:
                        pass

    # 2. Uploaded documents
    if os.path.exists(UPLOADS_DIR):
        if log_fn:
            log_fn("Scanning: " + UPLOADS_DIR)
        for f in os.listdir(UPLOADS_DIR):
            path = os.path.join(UPLOADS_DIR, f)
            if not os.path.isfile(path):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext not in (
                    ".pdf", ".pptx", ".docx",
                    ".doc", ".md", ".txt"):
                continue
            try:
                text = extract_text(path)
                if text.strip():
                    chunks = _chunk_text(
                        text, f, path, "uploaded")
                    docs.extend(
                        c[0] for c in chunks)
                    meta.extend(
                        c[1] for c in chunks)
                    if log_fn:
                        log_fn(
                            "  + " + f + " ("
                            + str(len(chunks))
                            + " chunks)")
            except Exception as exc:
                if log_fn:
                    log_fn(
                        "  ! Error: " + f
                        + " - " + str(exc))

    if not docs:
        if log_fn:
            log_fn("No content found!")
        return None, []

    if log_fn:
        log_fn("Encoding " + str(len(docs))
               + " chunks...")
    emb = embed_model.encode(docs).astype(
        "float32")
    idx = faiss.IndexFlatL2(emb.shape[1])
    idx.add(emb)
    faiss.write_index(idx, INDEX_FILE)
    pickle.dump(meta, open(META_FILE, "wb"))
    if log_fn:
        log_fn("Index built: "
               + str(len(docs)) + " chunks.")
    index = idx
    metadata = meta
    return idx, meta


def add_document_to_index(filepath, log_fn=None):
    global index, metadata
    if embed_model is None:
        if log_fn:
            log_fn("Vault not loaded yet!")
        return False
    fname = os.path.basename(filepath)
    if log_fn:
        log_fn("Extracting: " + fname)
    text = extract_text(filepath)
    if not text.strip():
        if log_fn:
            log_fn("No text in: " + fname)
        return False
    chunks = _chunk_text(
        text, fname, filepath, "uploaded")
    if not chunks:
        if log_fn:
            log_fn("No chunks from: " + fname)
        return False
    if log_fn:
        log_fn("Encoding " + str(len(chunks))
               + " chunks...")
    new_docs = [c[0] for c in chunks]
    new_meta = [c[1] for c in chunks]
    emb = embed_model.encode(new_docs).astype(
        "float32")
    if index is None:
        index = faiss.IndexFlatL2(emb.shape[1])
        metadata = []
    index.add(emb)
    metadata.extend(new_meta)
    faiss.write_index(index, INDEX_FILE)
    pickle.dump(metadata, open(META_FILE, "wb"))
    if log_fn:
        log_fn("Added " + str(len(chunks))
               + " chunks from " + fname
               + " (total: "
               + str(index.ntotal) + ")")
    return True


def retrieve(query):
    q = embed_model.encode([query]).astype(
        "float32")
    _, ids = index.search(q, TOP_K)
    results = []
    for i in ids[0]:
        if 0 <= i < len(metadata):
            results.append(metadata[i]["text"])
    return "\n\n".join(results)


def retrieve_with_sources(query):
    q = embed_model.encode([query]).astype(
        "float32")
    _, ids = index.search(q, TOP_K)
    results = []
    for i in ids[0]:
        if 0 <= i < len(metadata):
            m = metadata[i]
            results.append({
                "text": m["text"],
                "file": m.get("file", "?"),
                "source": m.get(
                    "source", "unknown")
            })
    return results


def build_prompt(context, topic):
        return [
            {"role": "system", "content":
                "You MUST answer fully based on "
                "the provided context.\n\n"
                "Only include explanation, "
                "algorithms, pseudocode, or C code "
                "if those details are explicitly "
                "present in the context. "
                "If something is missing, say "
                "'Not available in the context.' "
                "Do NOT invent new code."},
        {"role": "user",
         "content": context
                    + "\n\nTopic/Question:\n"
                    + topic}
    ]


def stream_generate(messages, out_fn):
    stream = ollama.chat(
        model=MODEL, messages=messages,
        stream=True,
        options={"num_predict": MAX_TOKENS})
    full = ""
    for chunk in stream:
        t = chunk["message"]["content"]
        out_fn(t)
        full += t
    return full


def append_log(q, a):
    logs = []
    if os.path.exists(LOG_FILE):
        try:
            logs = json.load(open(
                LOG_FILE, "r", encoding="utf-8"))
        except Exception:
            pass
    logs.append({
        "timestamp": datetime.now().isoformat(),
        "query": q, "answer": a})
    json.dump(logs, open(
        LOG_FILE, "w", encoding="utf-8"), indent=2)


# ==================================================
#  CHROMIUM BROWSER SCRIPT
# ==================================================

_DL_DIR_ESCAPED = DOWNLOAD_DIR.replace(
    "\\", "\\\\")

BROWSER_SCRIPT = r'''
import sys
import os

start_url = sys.argv[1] if len(sys.argv) > 1 else "__HOME__"
home_url = "__HOME__"
dl_dir = "__DL_DIR__"

try:
    from PySide6.QtWidgets import (
        QApplication, QMainWindow, QToolBar,
        QLineEdit, QSizePolicy, QMessageBox,
        QProgressBar)
    from PySide6.QtWebEngineWidgets import (
        QWebEngineView)
    from PySide6.QtWebEngineCore import (
        QWebEnginePage, QWebEngineProfile,
        QWebEngineSettings)
    from PySide6.QtCore import QUrl, Qt, QTimer
    from PySide6.QtGui import QAction
except ImportError:
    try:
        from PyQt6.QtWidgets import (
            QApplication, QMainWindow, QToolBar,
            QLineEdit, QSizePolicy, QMessageBox,
            QProgressBar)
        from PyQt6.QtWebEngineWidgets import (
            QWebEngineView)
        from PyQt6.QtWebEngineCore import (
            QWebEnginePage, QWebEngineProfile,
            QWebEngineSettings)
        from PyQt6.QtCore import QUrl, Qt, QTimer
        from PyQt6.QtGui import QAction
    except ImportError:
        print("NO_ENGINE", file=sys.stderr)
        sys.exit(1)


class InternalPage(QWebEnginePage):
    def __init__(self, profile, parent=None):
        super().__init__(profile, parent)
        self._temp_pages = []

    def createWindow(self, window_type):
        temp = QWebEnginePage(self.profile(), self)
        self._temp_pages.append(temp)
        temp.urlChanged.connect(
            lambda url, t=temp:
                self._on_popup_url(url, t))
        temp.loadFinished.connect(
            lambda ok, t=temp:
                self._cleanup_temp(t))
        return temp

    def _on_popup_url(self, url, temp_page):
        url_str = url.toString()
        if (url_str
                and url_str != "about:blank"
                and not url_str.startswith("blob:")):
            self.setUrl(url)

    def _cleanup_temp(self, temp_page):
        QTimer.singleShot(
            3000, lambda t=temp_page:
                self._remove_temp(t))

    def _remove_temp(self, temp_page):
        if temp_page in self._temp_pages:
            self._temp_pages.remove(temp_page)
            try:
                temp_page.deleteLater()
            except Exception:
                pass


class Browser(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("TheHelper")
        self.resize(1100, 750)
        self.setMinimumSize(800, 500)
        self._downloads = []

        self.setStyleSheet("""
            QMainWindow{background:#1a1a1a}
            QToolBar{background:#222;border:none;
                spacing:4px;padding:6px}
            QToolBar QToolButton{color:#ccc;
                padding:4px 8px;border:none;
                border-radius:3px;font-size:12px}
            QToolBar QToolButton:hover{background:#444}
            QLineEdit{background:#333;color:#eee;
                border:1px solid #444;
                border-radius:4px;
                padding:5px 10px;font-size:13px}
            QLineEdit:focus{border:1px solid #e8c547}
            QStatusBar{background:#1a1a1a;
                color:#aaa;font-size:12px}
            QProgressBar{background:#333;
                border:1px solid #444;
                border-radius:3px;text-align:center;
                color:#eee;font-size:11px;
                max-height:16px}
            QProgressBar::chunk{background:#e8c547;
                border-radius:2px}
        """)

        self.profile = QWebEngineProfile.defaultProfile()
        os.makedirs(dl_dir, exist_ok=True)
        try:
            self.profile.setDownloadPath(dl_dir)
        except Exception:
            pass
        self.profile.downloadRequested.connect(
            self.handle_download)

        self.view = QWebEngineView()
        self.page = InternalPage(
            self.profile, self.view)
        self.view.setPage(self.page)

        settings = self.page.settings()
        try:
            settings.setAttribute(
                QWebEngineSettings.WebAttribute.JavascriptEnabled, True)
            settings.setAttribute(
                QWebEngineSettings.WebAttribute.LocalStorageEnabled, True)
            settings.setAttribute(
                QWebEngineSettings.WebAttribute.PluginsEnabled, True)
        except Exception:
            pass

        nav = QToolBar()
        nav.setMovable(False)
        self.addToolBar(nav)

        for label, slot in [
                ("Back", self.view.back),
                ("Fwd", self.view.forward),
                ("Reload", self.view.reload),
                ("Home", self.go_home)]:
            a = QAction(label, self)
            a.triggered.connect(slot)
            nav.addAction(a)

        nav.addSeparator()
        self.url_bar = QLineEdit()
        self.url_bar.setPlaceholderText("Enter URL...")
        sp = QSizePolicy(
            QSizePolicy.Policy.Expanding,
            QSizePolicy.Policy.Preferred)
        self.url_bar.setSizePolicy(sp)
        self.url_bar.returnPressed.connect(self.navigate)
        nav.addWidget(self.url_bar)
        go = QAction("Go", self)
        go.triggered.connect(self.navigate)
        nav.addAction(go)
        nav.addSeparator()
        dl_open = QAction("Downloads", self)
        dl_open.triggered.connect(self.open_dl_folder)
        nav.addAction(dl_open)

        self.progress = QProgressBar()
        self.progress.setFixedWidth(200)
        self.progress.setVisible(False)
        self.statusBar().addPermanentWidget(self.progress)

        self.view.urlChanged.connect(self.on_url_changed)
        self.view.titleChanged.connect(self.on_title_changed)
        self.view.loadStarted.connect(
            lambda: self.statusBar().showMessage("Loading..."))
        self.view.loadFinished.connect(self.on_load_finished)

        self.setCentralWidget(self.view)
        self.statusBar().showMessage("Starting...")
        self.view.setUrl(QUrl(start_url))

    def navigate(self):
        u = self.url_bar.text().strip()
        if not u:
            return
        if not u.startswith("http"):
            u = "https://" + u
        self.view.setUrl(QUrl(u))

    def go_home(self):
        self.view.setUrl(QUrl(home_url))

    def on_url_changed(self, qurl):
        self.url_bar.setText(qurl.toString())

    def on_title_changed(self, title):
        self.setWindowTitle(
            (title + " - TheHelper") if title else "TheHelper")

    def on_load_finished(self, ok):
        self.statusBar().showMessage(
            "Ready" if ok else "Page load error")

    def open_dl_folder(self):
        p = os.path.realpath(dl_dir)
        os.makedirs(p, exist_ok=True)
        if sys.platform == "win32":
            os.startfile(p)
        elif sys.platform == "darwin":
            os.system('open "' + p + '"')
        else:
            os.system('xdg-open "' + p + '"')

    def handle_download(self, download):
        name = self._get_dl_name(download)
        os.makedirs(dl_dir, exist_ok=True)
        self._set_dl_path(download, name)
        download.accept()
        self._downloads.append(download)
        self.statusBar().showMessage("Downloading: " + name + "...")
        self.progress.setValue(0)
        self.progress.setVisible(True)
        self.progress.setFormat(name + " %p%")
        try:
            download.receivedBytesChanged.connect(
                lambda _n=name, _d=download:
                    self._on_progress(_d, _n))
        except Exception:
            pass
        try:
            download.isFinishedChanged.connect(
                lambda _n=name, _d=download:
                    self._on_finished(_d, _n))
        except Exception:
            try:
                download.finished.connect(
                    lambda _n=name, _d=download:
                        self._on_finished(_d, _n))
            except Exception:
                pass

    def _get_dl_name(self, download):
        name = ""
        for method in ["downloadFileName",
                "suggestedFileName", "fileName"]:
            try:
                name = getattr(download, method)()
                if name:
                    break
            except Exception:
                continue
        if not name:
            try:
                url = download.url().toString()
                name = url.split("/")[-1].split("?")[0]
            except Exception:
                pass
        return name or "download"

    def _set_dl_path(self, download, name):
        try:
            download.setDownloadDirectory(dl_dir)
            download.setDownloadFileName(name)
            return
        except Exception:
            pass
        try:
            download.setPath(os.path.join(dl_dir, name))
        except Exception:
            pass

    def _on_progress(self, download, name):
        try:
            received = download.receivedBytes()
            total = download.totalBytes()
            if total > 0:
                pct = int(received * 100 / total)
                self.progress.setValue(pct)
                self.statusBar().showMessage(
                    "Downloading " + name + ": " + str(pct) + "%")
            else:
                self.progress.setRange(0, 0)
        except Exception:
            pass

    def _on_finished(self, download, name):
        if download in self._downloads:
            self._downloads.remove(download)
        self.progress.setRange(0, 100)
        self.progress.setValue(100)
        self.statusBar().showMessage(
            "Saved: " + name + " -> " + dl_dir)
        QTimer.singleShot(3000,
            lambda: self.progress.setVisible(False))
        msg = QMessageBox(self)
        msg.setWindowTitle("Download Complete")
        msg.setText("Saved:\n" + os.path.join(dl_dir, name))
        msg.setInformativeText("Open downloads folder?")
        msg.setStandardButtons(
            QMessageBox.StandardButton.Yes
            | QMessageBox.StandardButton.No)
        result = msg.exec()
        if result == QMessageBox.StandardButton.Yes:
            self.open_dl_folder()


app = QApplication(sys.argv)
win = Browser()
win.show()
sys.exit(app.exec())
'''.replace(
    "__HOME__", BASE_URL
).replace(
    "__DL_DIR__", _DL_DIR_ESCAPED
)


# ==================================================
#  GUI
# ==================================================

class VaultGUI:
    def __init__(self, root):
        self.root = root
        self.root.title(
            "Vault AI + TheHelper Browser")
        self.root.geometry("1260x780")
        self.root.minsize(900, 600)
        self.root.configure(bg=BG)

        self._vault_ready = False
        self._browser_proc = None
        self._browser_launched = False

        self._apply_styles()
        self._build_ui()

        threading.Thread(
            target=self._init_vault,
            daemon=True).start()

    def _apply_styles(self):
        s = ttk.Style()
        s.theme_use("clam")
        s.configure(
            ".", background=BG, foreground=FG,
            fieldbackground=BG2,
            bordercolor=BORDER,
            troughcolor=BG3, relief="flat")
        s.configure(
            "TNotebook", background=BG,
            borderwidth=0)
        s.configure(
            "TNotebook.Tab", background=BG2,
            foreground=FG_DIM, padding=[16, 6],
            font=FONT_UI, borderwidth=0)
        s.map("TNotebook.Tab",
              background=[("selected", BG3)],
              foreground=[("selected", ACCENT)])
        s.configure("TFrame", background=BG)
        s.configure(
            "TLabel", background=BG,
            foreground=FG, font=FONT_UI)
        s.configure(
            "TButton", background=BG3,
            foreground=FG, font=FONT_UI,
            relief="flat", borderwidth=0,
            padding=[10, 5])
        s.map("TButton",
              background=[("active", ACCENT),
                          ("pressed", ACCENT)],
              foreground=[("active", BG)])
        s.configure(
            "Accent.TButton",
            background=ACCENT, foreground=BG,
            font=(*FONT_UI[:2], "bold"),
            relief="flat", padding=[10, 5])
        s.map("Accent.TButton",
              background=[("active", "#f5d76e")])
        s.configure(
            "TEntry", fieldbackground=BG2,
            foreground=FG, insertcolor=ACCENT,
            relief="flat", padding=6)

    def _build_ui(self):
        header = tk.Frame(
            self.root, bg=BG, pady=10)
        header.pack(fill="x", padx=20)
        tk.Label(
            header, text="VAULT AI",
            font=("Courier New", 18, "bold"),
            bg=BG, fg=ACCENT).pack(side="left")
        tk.Label(
            header,
            text=" + TheHelper Browser",
            font=("Segoe UI", 11), bg=BG,
            fg=FG_DIM).pack(
                side="left", padx=4)
        self._status_var = tk.StringVar(
            value="Initialising...")
        tk.Label(
            header,
            textvariable=self._status_var,
            font=FONT_UI, bg=BG,
            fg=FG_DIM).pack(side="right")

        tk.Frame(
            self.root, bg=BORDER,
            height=1).pack(fill="x")

        self.nb = ttk.Notebook(self.root)
        self.nb.pack(fill="both", expand=True)

        self._build_chat_tab()      # Tab 0
        self._build_docs_tab()      # Tab 1
        self._build_helper_tab()    # Tab 2
        self._build_log_tab()       # Tab 3

        self.nb.bind(
            "<<NotebookTabChanged>>",
            self._on_tab_changed)
        self.root.protocol(
            "WM_DELETE_WINDOW", self._on_close)

    def _on_close(self):
        self._kill_browser()
        self.root.destroy()

    # ==================================================
    #  TAB 0 -- Vault AI Chat (with Upload button)
    # ==================================================
    def _build_chat_tab(self):
        tab = ttk.Frame(self.nb)
        self.nb.add(tab, text="  Vault AI  ")

        out_frame = tk.Frame(tab, bg=BG)
        out_frame.pack(
            fill="both", expand=True,
            padx=16, pady=(12, 4))

        self.chat_out = scrolledtext.ScrolledText(
            out_frame, wrap="word", bg=BG2,
            fg=FG,
            insertbackground=ACCENT,
            font=FONT_MONO,
            relief="flat", borderwidth=0,
            state="disabled",
            selectbackground=ACCENT,
            selectforeground=BG)
        self.chat_out.pack(
            fill="both", expand=True)
        for tag, color in [
                ("accent", ACCENT),
                ("dim", FG_DIM),
                ("green", GREEN),
                ("red", RED)]:
            self.chat_out.tag_configure(
                tag, foreground=color)
        self.chat_out.tag_configure(
            "heading", foreground=ACCENT2,
            font=(*FONT_MONO[:1],
                  FONT_MONO[1] + 1, "bold"))
        self.chat_out.tag_configure(
            "source", foreground="#9e9e9e",
            font=("Segoe UI", 8, "italic"))

        inp_frame = tk.Frame(
            tab, bg=BG, pady=8)
        inp_frame.pack(fill="x", padx=16)

        self.chat_input = tk.Text(
            inp_frame, height=4, bg=BG2,
            fg=FG,
            insertbackground=ACCENT,
            font=FONT_MONO,
            relief="flat", wrap="word",
            selectbackground=ACCENT,
            selectforeground=BG)
        self.chat_input.pack(
            side="left", fill="both",
            expand=True)
        self.chat_input.bind(
            "<Return>", self._on_enter)
        self.chat_input.bind(
            "<Control-j>",
            self._insert_newline)

        btn_col = tk.Frame(
            inp_frame, bg=BG, padx=6)
        btn_col.pack(side="left", fill="y")

        ttk.Button(
            btn_col, text="Send\n(Enter)",
            style="Accent.TButton",
            command=self._send_query).pack(
                fill="x", pady=(0, 4))

        ttk.Button(
            btn_col, text="Upload\nFiles",
            command=self._upload_from_chat
            ).pack(fill="x", pady=(0, 4))

        ttk.Button(
            btn_col, text="Clear",
            command=self._clear_chat).pack(
                fill="x")

        info_frame = tk.Frame(tab, bg=BG)
        info_frame.pack(fill="x", padx=16)
        tk.Label(
            info_frame,
            text="Enter = Send | "
                 "Ctrl+J = Newline | "
                 "Upload PDF/PPTX/DOCX",
            font=("Segoe UI", 8), bg=BG,
            fg=FG_DIM).pack(
                side="left", pady=2)

        self._index_info = tk.StringVar(value="")
        tk.Label(
            info_frame,
            textvariable=self._index_info,
            font=("Segoe UI", 8), bg=BG,
            fg=FG_DIM).pack(
                side="right", pady=2)

    def _on_enter(self, event):
        self._send_query()
        return "break"

    def _insert_newline(self, event):
        self.chat_input.insert("insert", "\n")
        return "break"

    def _send_query(self):
        q = self.chat_input.get(
            "1.0", "end").strip()
        if not q:
            return
        self.chat_input.delete("1.0", "end")
        threading.Thread(
            target=self._run_query,
            args=(q,), daemon=True).start()

    def _run_query(self, query):
        if not self._vault_ready:
            self._chat_write(
                "Vault not ready yet.\n", "red")
            return
        if index is None or index.ntotal == 0:
            self._chat_write(
                "No documents indexed!\n"
                "Upload docs in the Documents "
                "tab first.\n", "red")
            return

        raw_lines = query.splitlines()
        bullet_topics = []
        for line in raw_lines:
            stripped = line.strip()
            if not stripped:
                continue
            if stripped[0] in ("-", "*"):
                candidate = stripped[1:].strip()
                if candidate:
                    bullet_topics.append(candidate)
        if bullet_topics:
            topics = bullet_topics
        else:
            topics = [query.strip()]
        full = ""
        for t in topics:
            self._chat_write(
                "\n== " + t + " ==\n",
                "heading")
            results = retrieve_with_sources(t)
            if not results:
                self._chat_write(
                    "\nNo indexed content found for:\n"
                    + t + "\n"
                    "Upload matching documents first.\n"
                    "-- Skipped --\n", "red")
                continue
            sources = set()
            for r in results:
                sources.add(r["file"])
            if sources:
                self._chat_write(
                    "Sources: "
                    + ", ".join(sources)
                    + "\n", "source")
            ctx = "\n\n".join(
                r["text"] for r in results)
            ans = stream_generate(
                build_prompt(ctx, t),
                lambda x: self._chat_write(x))
            full += "\n\n== " + t + " ==\n" + ans
        append_log(query, full)
        self._chat_write(
            "\n\n-- Done --\n", "dim")

    def _chat_write(self, text, tag=None):
        def _do():
            self.chat_out.configure(
                state="normal")
            if tag:
                self.chat_out.insert(
                    "end", text, (tag,))
            else:
                self.chat_out.insert(
                    "end", text)
            self.chat_out.see("end")
            self.chat_out.configure(
                state="disabled")
        self.root.after(0, _do)

    def _clear_chat(self):
        self.chat_out.configure(state="normal")
        self.chat_out.delete("1.0", "end")
        self.chat_out.configure(
            state="disabled")

    def _update_index_info(self):
        if index is not None:
            count = str(index.ntotal)
            ndocs = str(len(load_docs_log()))
            self._index_info.set(
                "Index: " + count
                + " chunks | "
                + ndocs + " uploaded docs")
        else:
            self._index_info.set("Index: empty")

    # -- Upload from chat tab --
    def _upload_from_chat(self):
        filetypes = get_supported_types()
        paths = filedialog.askopenfilenames(
            title="Select documents to upload",
            filetypes=filetypes)
        if not paths:
            return
        self._chat_write(
            "\nUploading " + str(len(paths))
            + " file(s)...\n", "accent")

        def _do_upload():
            os.makedirs(
                UPLOADS_DIR, exist_ok=True)
            docs_log = load_docs_log()
            success = 0
            failed = 0

            for path in paths:
                fname = os.path.basename(path)
                self._chat_write(
                    "  Processing: " + fname
                    + "\n", "dim")

                dest = os.path.join(
                    UPLOADS_DIR, fname)
                try:
                    shutil.copy2(path, dest)
                except Exception as exc:
                    err = str(exc)
                    self._chat_write(
                        "  Copy failed: "
                        + err + "\n", "red")
                    failed += 1
                    continue

                ext = os.path.splitext(
                    fname)[1].lower()
                if ext == ".pdf" and not HAS_PDF:
                    self._chat_write(
                        "  Skipped (install "
                        "PyPDF2)\n", "red")
                    failed += 1
                    continue
                if (ext == ".pptx"
                        and not HAS_PPTX):
                    self._chat_write(
                        "  Skipped (install "
                        "python-pptx)\n", "red")
                    failed += 1
                    continue
                if (ext in (".docx", ".doc")
                        and not HAS_DOCX):
                    self._chat_write(
                        "  Skipped (install "
                        "python-docx)\n", "red")
                    failed += 1
                    continue

                text = extract_text(dest)
                chars = len(text.strip())
                if chars == 0:
                    self._chat_write(
                        "  No text in "
                        + fname + "\n", "red")
                    failed += 1
                    continue

                self._chat_write(
                    "  Extracted "
                    + str(chars)
                    + " chars\n", "dim")

                if self._vault_ready:
                    ok = add_document_to_index(
                        dest,
                        lambda m:
                            self._chat_write(
                                "  " + m
                                + "\n", "dim"))
                    if ok:
                        success += 1
                    else:
                        failed += 1
                else:
                    success += 1

                existing = [
                    d for d in docs_log
                    if d["name"] != fname]
                existing.append({
                    "name": fname,
                    "path": dest,
                    "ext": ext,
                    "chars": chars,
                    "uploaded": (
                        datetime.now()
                        .isoformat()),
                    "indexed": self._vault_ready
                })
                docs_log = existing
                save_docs_log(docs_log)

            self._chat_write(
                "\nUpload complete: "
                + str(success) + " indexed, "
                + str(failed) + " failed\n",
                "green" if failed == 0
                else "accent")
            if success > 0:
                self._chat_write(
                    "Ask questions about your "
                    "documents now!\n", "green")
            self.root.after(
                0, self._refresh_docs_list)
            self.root.after(
                0, self._update_index_info)

        threading.Thread(
            target=_do_upload,
            daemon=True).start()

    # ==================================================
    #  TAB 1 -- Documents (manage uploads)
    # ==================================================
    def _build_docs_tab(self):
        tab = ttk.Frame(self.nb)
        self.nb.add(tab, text="  Documents  ")

        # Top controls
        top = tk.Frame(tab, bg=BG, pady=8)
        top.pack(fill="x", padx=16)

        ttk.Button(
            top, text="Upload Files",
            style="Accent.TButton",
            command=self._upload_files_from_docs
            ).pack(side="left", padx=(0, 8))
        ttk.Button(
            top, text="Upload Folder",
            command=self._upload_folder
            ).pack(side="left", padx=(0, 8))
        ttk.Button(
            top, text="Rebuild Index",
            command=self._rebuild_index
            ).pack(side="left", padx=(0, 8))
        ttk.Button(
            top, text="Open Uploads Folder",
            command=lambda: (
                os.makedirs(
                    UPLOADS_DIR, exist_ok=True),
                os.startfile(UPLOADS_DIR)
                if sys.platform == "win32"
                else webbrowser.open(UPLOADS_DIR))
            ).pack(side="right", padx=4)

        # Status
        self._docs_status = tk.StringVar(
            value="  Ready to upload")
        tk.Label(
            tab,
            textvariable=self._docs_status,
            bg=BG, fg=FG_DIM,
            font=("Segoe UI", 8),
            anchor="w").pack(
                fill="x", padx=16)

        # Dep status
        dep_frame = tk.Frame(
            tab, bg=BG3, padx=12, pady=8)
        dep_frame.pack(
            fill="x", padx=16, pady=(4, 8))
        tk.Label(
            dep_frame, text="Parser Status:",
            font=(*FONT_UI[:2], "bold"),
            bg=BG3, fg=ACCENT,
            anchor="w").pack(fill="x")
        tk.Label(
            dep_frame, text=get_dep_status(),
            font=("Consolas", 9), bg=BG3,
            fg=GREEN if (
                HAS_PDF and HAS_PPTX
                and HAS_DOCX) else FG_DIM,
            anchor="w",
            justify="left").pack(fill="x")

        # Document list + preview
        cols = tk.Frame(tab, bg=BG)
        cols.pack(
            fill="both", expand=True,
            padx=16, pady=(0, 12))

        list_frame = tk.Frame(cols, bg=BG)
        list_frame.pack(
            side="left", fill="both",
            expand=True, padx=(0, 8))

        tk.Label(
            list_frame,
            text="Uploaded Documents",
            font=(*FONT_UI[:2], "bold"),
            bg=BG, fg=FG,
            anchor="w").pack(fill="x")

        list_container = tk.Frame(
            list_frame, bg=BG2)
        list_container.pack(
            fill="both", expand=True,
            pady=(4, 0))

        self.docs_listbox = tk.Listbox(
            list_container, bg=BG2, fg=FG,
            selectbackground=ACCENT,
            selectforeground=BG,
            font=FONT_MONO, relief="flat",
            borderwidth=0,
            activestyle="none")
        scrollbar = ttk.Scrollbar(
            list_container, orient="vertical",
            command=self.docs_listbox.yview)
        self.docs_listbox.configure(
            yscrollcommand=scrollbar.set)
        self.docs_listbox.pack(
            side="left", fill="both",
            expand=True)
        scrollbar.pack(
            side="right", fill="y")

        list_btns = tk.Frame(
            list_frame, bg=BG, pady=4)
        list_btns.pack(fill="x")
        ttk.Button(
            list_btns, text="Remove Selected",
            command=self._remove_selected
            ).pack(side="left", padx=(0, 4))
        ttk.Button(
            list_btns, text="Remove All",
            command=self._remove_all
            ).pack(side="left", padx=(0, 4))
        ttk.Button(
            list_btns, text="Refresh",
            command=self._refresh_docs_list
            ).pack(side="left")

        info_frame = tk.Frame(
            cols, bg=BG, width=300)
        info_frame.pack(
            side="right", fill="both",
            padx=(8, 0))
        info_frame.pack_propagate(False)

        tk.Label(
            info_frame, text="Preview",
            font=(*FONT_UI[:2], "bold"),
            bg=BG, fg=FG,
            anchor="w").pack(fill="x")

        self.doc_preview = (
            scrolledtext.ScrolledText(
                info_frame, wrap="word",
                bg=BG2, fg=FG,
                font=("Consolas", 9),
                relief="flat",
                state="disabled",
                height=15))
        self.doc_preview.pack(
            fill="both", expand=True,
            pady=(4, 0))

        self.docs_listbox.bind(
            "<<ListboxSelect>>",
            self._on_doc_selected)

        self._refresh_docs_list()

    def _upload_files_from_docs(self):
        filetypes = get_supported_types()
        paths = filedialog.askopenfilenames(
            title="Select documents",
            filetypes=filetypes)
        if not paths:
            return
        threading.Thread(
            target=self._process_uploads,
            args=(list(paths),),
            daemon=True).start()

    def _upload_folder(self):
        folder = filedialog.askdirectory(
            title="Select folder")
        if not folder:
            return
        paths = []
        for f in os.listdir(folder):
            ext = os.path.splitext(f)[1].lower()
            if ext in (".pdf", ".pptx", ".docx",
                       ".doc", ".md", ".txt"):
                paths.append(
                    os.path.join(folder, f))
        if not paths:
            messagebox.showinfo(
                "No files",
                "No supported files found.")
            return
        threading.Thread(
            target=self._process_uploads,
            args=(paths,),
            daemon=True).start()

    def _process_uploads(self, paths):
        os.makedirs(UPLOADS_DIR, exist_ok=True)
        docs_log = load_docs_log()
        total = len(paths)
        success = 0
        failed = 0

        for i, path in enumerate(paths, 1):
            fname = os.path.basename(path)
            self._set_docs_status(
                "  Processing " + str(i)
                + "/" + str(total) + ": "
                + fname)

            dest = os.path.join(
                UPLOADS_DIR, fname)
            try:
                shutil.copy2(path, dest)
            except Exception as exc:
                self._set_docs_status(
                    "  Copy failed: "
                    + str(exc))
                failed += 1
                continue

            ext = os.path.splitext(
                fname)[1].lower()
            if ext == ".pdf" and not HAS_PDF:
                failed += 1
                continue
            if ext == ".pptx" and not HAS_PPTX:
                failed += 1
                continue
            if (ext in (".docx", ".doc")
                    and not HAS_DOCX):
                failed += 1
                continue

            text = extract_text(dest)
            chars = len(text.strip())
            if chars == 0:
                failed += 1
                continue

            if self._vault_ready:
                ok = add_document_to_index(
                    dest,
                    lambda m:
                        self._set_docs_status(
                            "  " + m))
                if ok:
                    success += 1
                else:
                    failed += 1
            else:
                success += 1

            existing = [
                d for d in docs_log
                if d["name"] != fname]
            existing.append({
                "name": fname,
                "path": dest,
                "ext": ext,
                "chars": chars,
                "uploaded": (
                    datetime.now().isoformat()),
                "indexed": self._vault_ready
            })
            docs_log = existing
            save_docs_log(docs_log)

        msg = (str(success) + " ok, "
               + str(failed) + " failed")
        self._set_docs_status(
            "  Upload done: " + msg)
        self.root.after(
            0, self._refresh_docs_list)
        self.root.after(
            0, self._update_index_info)
        if success > 0:
            self.root.after(0, lambda:
                messagebox.showinfo(
                    "Done",
                    str(success)
                    + " documents uploaded!"))

    def _refresh_docs_list(self):
        self.docs_listbox.delete(0, "end")
        docs = load_docs_log()
        if not docs:
            self.docs_listbox.insert(
                "end",
                "  No documents uploaded yet")
            self.docs_listbox.insert(
                "end", "")
            self.docs_listbox.insert(
                "end",
                "  Click 'Upload Files' or use")
            self.docs_listbox.insert(
                "end",
                "  the Upload button in Chat tab")
            return
        for doc in docs:
            name = doc["name"]
            chars = doc.get("chars", 0)
            ext = doc.get("ext", "")
            icon = {
                ".pdf": "[PDF]",
                ".pptx": "[PPT]",
                ".docx": "[DOC]",
                ".md": "[MD]",
                ".txt": "[TXT]"
            }.get(ext, "[???]")
            if chars > 10000:
                size = (" ("
                    + str(chars // 1000)
                    + "k chars)")
            elif chars > 0:
                size = (" ("
                    + str(chars) + " chars)")
            else:
                size = ""
            self.docs_listbox.insert(
                "end",
                "  " + icon + "  "
                + name + size)

    def _on_doc_selected(self, event):
        sel = self.docs_listbox.curselection()
        if not sel:
            return
        docs = load_docs_log()
        idx = sel[0]
        if idx >= len(docs):
            return
        doc = docs[idx]
        path = doc.get("path", "")

        self.doc_preview.configure(
            state="normal")
        self.doc_preview.delete("1.0", "end")
        info = (
            "File: " + doc["name"] + "\n"
            "Type: " + doc.get("ext", "?")
            + "\n"
            "Chars: " + str(
                doc.get("chars", 0)) + "\n"
            "Uploaded: " + doc.get(
                "uploaded", "?")[:19] + "\n"
            + "-" * 40 + "\n\n")
        self.doc_preview.insert("end", info)
        if os.path.exists(path):
            text = extract_text(path)
            preview = text[:2000]
            if len(text) > 2000:
                preview += ("\n\n... ("
                    + str(len(text))
                    + " total chars)")
            self.doc_preview.insert(
                "end", preview)
        else:
            self.doc_preview.insert(
                "end", "(File not found)")
        self.doc_preview.configure(
            state="disabled")

    def _remove_selected(self):
        sel = self.docs_listbox.curselection()
        if not sel:
            return
        docs = load_docs_log()
        idx = sel[0]
        if idx >= len(docs):
            return
        doc = docs[idx]
        name = doc["name"]
        if not messagebox.askyesno(
                "Remove",
                "Remove " + name + "?"):
            return
        path = doc.get("path", "")
        if os.path.exists(path):
            try:
                os.remove(path)
            except Exception:
                pass
        docs.pop(idx)
        save_docs_log(docs)
        self._refresh_docs_list()
        self._rebuild_index()

    def _remove_all(self):
        docs = load_docs_log()
        if not docs:
            return
        if not messagebox.askyesno(
                "Remove All",
                "Remove all " + str(len(docs))
                + " documents?"):
            return
        for doc in docs:
            path = doc.get("path", "")
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass
        save_docs_log([])
        self._refresh_docs_list()
        self._rebuild_index()

    def _rebuild_index(self):
        if not self._vault_ready:
            messagebox.showinfo(
                "Not ready",
                "Vault is still loading.")
            return
        self._set_docs_status(
            "  Rebuilding index...")
        threading.Thread(
            target=self._do_rebuild,
            daemon=True).start()

    def _do_rebuild(self):
        for f in [INDEX_FILE, META_FILE]:
            if os.path.exists(f):
                try:
                    os.remove(f)
                except Exception:
                    pass
        rebuild_full_index(
            lambda m: self._set_docs_status(
                "  " + m))
        self.root.after(
            0, self._update_index_info)
        self._set_docs_status(
            "  Index rebuilt!")

    def _set_docs_status(self, msg):
        self.root.after(
            0,
            lambda: self._docs_status.set(msg))

    # ==================================================
    #  TAB 2 -- TheHelper Browser
    # ==================================================
    def _build_helper_tab(self):
        self._helper_tab = ttk.Frame(self.nb)
        self.nb.add(
            self._helper_tab,
            text="  TheHelper  ")

        toolbar = tk.Frame(
            self._helper_tab,
            bg=BG2, pady=6, padx=8)
        toolbar.pack(fill="x")

        self._url_var = tk.StringVar(
            value=BASE_URL)
        url_entry = ttk.Entry(
            toolbar,
            textvariable=self._url_var,
            font=FONT_UI)
        url_entry.pack(
            side="left", fill="x",
            expand=True, padx=(0, 8))
        url_entry.bind(
            "<Return>",
            lambda e: self._open_url(
                self._url_var.get()))

        ttk.Button(
            toolbar, text="Go",
            command=lambda: self._open_url(
                self._url_var.get())
            ).pack(side="left", padx=(0, 8))
        ttk.Button(
            toolbar,
            text="Launch Browser",
            style="Accent.TButton",
            command=self._launch_browser
            ).pack(side="left", padx=(0, 8))
        ttk.Button(
            toolbar,
            text="Open Downloads",
            command=self._open_dl_folder
            ).pack(side="right", padx=4)
        ttk.Button(
            toolbar,
            text="System Browser",
            command=lambda: webbrowser.open(
                self._url_var.get())
            ).pack(side="right", padx=4)

        self.dl_dir_var = tk.StringVar(
            value=DOWNLOAD_DIR)

        self._browser_status = tk.StringVar(
            value="  Click tab or Launch")
        tk.Label(
            self._helper_tab,
            textvariable=self._browser_status,
            bg=BG, fg=FG_DIM,
            font=("Segoe UI", 8),
            anchor="w").pack(
                fill="x", padx=10)

        self._browser_content = tk.Frame(
            self._helper_tab, bg=BG)
        self._browser_content.pack(
            fill="both", expand=True)

        if HAS_QT_BROWSER:
            self._build_ready_screen()
        else:
            self._build_install_screen()

    def _open_dl_folder(self):
        os.makedirs(DOWNLOAD_DIR, exist_ok=True)
        if sys.platform == "win32":
            os.startfile(DOWNLOAD_DIR)
        else:
            webbrowser.open(DOWNLOAD_DIR)

    def _build_ready_screen(self):
        for w in (self._browser_content
                  .winfo_children()):
            w.destroy()
        f = tk.Frame(
            self._browser_content, bg=BG)
        f.pack(expand=True)
        tk.Label(
            f, text="TheHelper Browser",
            font=FONT_H1, bg=BG,
            fg=FG).pack(pady=(60, 5))
        tk.Label(
            f,
            text=(
                "Chromium browser "
                "(Qt WebEngine).\n"
                "All links stay inside.\n"
                "Downloads to: "
                + DOWNLOAD_DIR),
            font=FONT_UI, bg=BG, fg=FG_DIM,
            justify="center").pack(pady=20)
        ttk.Button(
            f,
            text="Launch TheHelper Now",
            style="Accent.TButton",
            command=self._launch_browser
            ).pack(pady=10)

    def _build_active_screen(self):
        for w in (self._browser_content
                  .winfo_children()):
            w.destroy()
        f = tk.Frame(
            self._browser_content, bg=BG)
        f.pack(expand=True)
        tk.Label(
            f,
            text="Chromium browser is running",
            font=FONT_H1, bg=BG,
            fg=GREEN).pack(pady=(80, 10))
        tk.Label(
            f, text="Alt+Tab to switch.",
            font=FONT_UI, bg=BG, fg=FG_DIM,
            justify="center").pack(pady=15)
        btn_row = tk.Frame(f, bg=BG)
        btn_row.pack(pady=10)
        ttk.Button(
            btn_row, text="Relaunch",
            style="Accent.TButton",
            command=self._relaunch_browser
            ).pack(side="left", padx=6)
        ttk.Button(
            btn_row, text="Open Downloads",
            command=self._open_dl_folder
            ).pack(side="left", padx=6)

    def _build_install_screen(self):
        for w in (self._browser_content
                  .winfo_children()):
            w.destroy()
        f = tk.Frame(
            self._browser_content, bg=BG)
        f.pack(expand=True)
        tk.Label(
            f, text="pip install PySide6",
            font=("Consolas", 14),
            bg=BG, fg=ACCENT).pack(pady=40)
        ttk.Button(
            f,
            text="Open in System Browser",
            style="Accent.TButton",
            command=lambda: webbrowser.open(
                BASE_URL)
            ).pack(pady=20)

    # Tab 2 is now index 2 (0=Chat, 1=Docs, 2=Helper)
    def _on_tab_changed(self, event):
        current = self.nb.index(
            self.nb.select())
        if (current == 2
                and HAS_QT_BROWSER
                and not self._browser_launched):
            self.root.after(
                300, self._launch_browser)

    # -- Browser process management --
    def _is_browser_running(self):
        if self._browser_proc is None:
            return False
        return self._browser_proc.poll() is None

    def _kill_browser(self):
        if self._browser_proc is not None:
            try:
                self._browser_proc.terminate()
                self._browser_proc.wait(
                    timeout=3)
            except Exception:
                try:
                    self._browser_proc.kill()
                except Exception:
                    pass
            self._browser_proc = None

    def _launch_browser(self):
        if self._is_browser_running():
            self._browser_status.set(
                "  Already running (Alt+Tab)")
            return
        if not HAS_QT_BROWSER:
            messagebox.showerror(
                "Missing",
                "pip install PySide6")
            return

        url = self._url_var.get() or BASE_URL
        self._browser_launched = True
        self._browser_status.set(
            "  Launching: " + url)
        self._build_active_screen()

        try:
            flags = 0
            if sys.platform == "win32":
                flags = (
                    subprocess.CREATE_NO_WINDOW)
            self._browser_proc = (
                subprocess.Popen(
                    [sys.executable,
                     "-c", BROWSER_SCRIPT,
                     url],
                    creationflags=flags,
                    stderr=subprocess.PIPE))
            self.root.after(
                1500,
                self._check_launch_success)
        except Exception as exc:
            err_msg = str(exc)
            self._browser_proc = None
            self._browser_status.set(
                "  Failed: " + err_msg)
            self._build_ready_screen()

    def _check_launch_success(self):
        if self._browser_proc is None:
            return
        ret = self._browser_proc.poll()
        if ret is not None:
            err = ""
            try:
                err = (
                    self._browser_proc
                    .stderr.read()
                    .decode("utf-8", "ignore")
                    .strip()[:300])
            except Exception:
                pass
            self._browser_proc = None
            self._browser_status.set(
                "  Crashed: "
                + (err[:80] or "unknown"))
            self._build_ready_screen()
            if err:
                messagebox.showerror(
                    "Error", err)
        else:
            pid = str(
                self._browser_proc.pid)
            self._browser_status.set(
                "  Running (PID: "
                + pid + ")")
            self._poll_browser()

    def _poll_browser(self):
        if self._is_browser_running():
            self.root.after(
                500, self._poll_browser)
        else:
            self._browser_proc = None
            self._browser_status.set(
                "  Closed -- Launch to reopen")
            self._build_ready_screen()

    def _relaunch_browser(self):
        self._kill_browser()
        self.root.after(
            300, self._launch_browser)

    def _open_url(self, url):
        if not url.startswith("http"):
            url = "https://" + url
        self._url_var.set(url)
        if self._is_browser_running():
            self._kill_browser()
            self.root.after(
                300, self._launch_browser)
        else:
            self._launch_browser()

    # ==================================================
    #  TAB 3 -- Chat Log
    # ==================================================
    def _build_log_tab(self):
        tab = ttk.Frame(self.nb)
        self.nb.add(tab, text="  Chat Log  ")

        ctrl = tk.Frame(tab, bg=BG, pady=8)
        ctrl.pack(fill="x", padx=16)
        ttk.Button(
            ctrl, text="Refresh",
            command=self._refresh_logs
            ).pack(side="left")
        ttk.Button(
            ctrl, text="Clear Log File",
            command=self._clear_log_file
            ).pack(side="left", padx=8)

        self.log_view = scrolledtext.ScrolledText(
            tab, wrap="word", bg=BG2, fg=FG,
            font=FONT_MONO, relief="flat",
            state="disabled",
            selectbackground=ACCENT,
            selectforeground=BG)
        self.log_view.pack(
            fill="both", expand=True,
            padx=16, pady=(0, 12))
        self.log_view.tag_configure(
            "q", foreground=ACCENT,
            font=(*FONT_MONO[:1],
                  FONT_MONO[1], "bold"))
        self.log_view.tag_configure(
            "ts", foreground=FG_DIM)
        self._refresh_logs()

    def _refresh_logs(self):
        self.log_view.configure(
            state="normal")
        self.log_view.delete("1.0", "end")
        if not os.path.exists(LOG_FILE):
            self.log_view.insert(
                "end", "No logs yet.")
        else:
            try:
                logs = json.load(open(
                    LOG_FILE, "r",
                    encoding="utf-8"))
                for entry in reversed(
                        logs[-50:]):
                    ts = entry[
                        "timestamp"][:19]
                    self.log_view.insert(
                        "end",
                        "[" + ts + "]\n",
                        "ts")
                    q_text = (
                        entry["query"][:120])
                    self.log_view.insert(
                        "end",
                        "Q: " + q_text + "\n",
                        "q")
                    a_text = (
                        entry["answer"][:400])
                    self.log_view.insert(
                        "end",
                        a_text + "\n\n")
            except Exception as exc:
                self.log_view.insert(
                    "end",
                    "Error: " + str(exc))
        self.log_view.configure(
            state="disabled")

    def _clear_log_file(self):
        if messagebox.askyesno(
                "Clear",
                "Delete all chat history?"):
            json.dump(
                [], open(LOG_FILE, "w"))
            self._refresh_logs()

    # -- Vault init --
    def _init_vault(self):
        self._set_status("Loading vault...")
        ok = load_vault(
            lambda m: self._chat_write(
                m + "\n", "dim"))
        if ok:
            self._vault_ready = True
            self._set_status("Vault ready")
            self._chat_write(
                "Vault AI ready.\n"
                "Upload documents (Upload "
                "button or Documents tab),\n"
                "then ask questions here.\n",
                "green")
            self.root.after(
                0, self._update_index_info)
        else:
            self._set_status(
                "Vault unavailable")
            self._chat_write(
                "Vault deps missing.\n"
                "pip install faiss-cpu "
                "sentence-transformers "
                "ollama\n", "red")

    def _set_status(self, msg):
        self.root.after(
            0,
            lambda: self._status_var.set(msg))


# ==================================================
if __name__ == "__main__":
    root = tk.Tk()
    app = VaultGUI(root)
    root.mainloop()
