#!/usr/bin/env python3
"""
VaultAI - Complete Standalone Application
Single file, no external vaultai package required.
Models dropdown shows ONLY locally installed Ollama models.
Split Upload button in Chat tab with vault-docs dropdown.
Compatible with both old and new ollama Python client APIs.
"""

import os
import sys
import json
import pickle
import shutil
import threading
import subprocess
import webbrowser
import importlib.util
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog, colorchooser
from datetime import datetime

# ============================================================
#  OPTIONAL DEPENDENCIES
# ============================================================

try:
    import faiss
    import numpy as np
    from sentence_transformers import SentenceTransformer
    import ollama
    VAULT_AVAILABLE = True
except ImportError:
    VAULT_AVAILABLE = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    HAS_QT_BROWSER = (
        importlib.util.find_spec("PySide6") is not None
        or importlib.util.find_spec("PyQt6") is not None
    )
except Exception:
    HAS_QT_BROWSER = False


# ============================================================
#  PATHS & CONFIG
# ============================================================

def _get_data_dir() -> str:
    if sys.platform == "win32":
        base = os.environ.get("LOCALAPPDATA", os.path.expanduser("~"))
    elif sys.platform == "darwin":
        base = os.path.join(
            os.path.expanduser("~"), "Library", "Application Support"
        )
    else:
        base = os.environ.get(
            "XDG_DATA_HOME",
            os.path.join(os.path.expanduser("~"), ".local", "share"),
        )
    path = os.path.join(base, "VaultAI")
    os.makedirs(path, exist_ok=True)
    return path


DATA_DIR      = _get_data_dir()
NOTES_FOLDER  = os.path.expanduser("~/Documents/Obsidian")
UPLOADS_DIR   = os.path.join(DATA_DIR, "uploaded_docs")
INDEX_FILE    = os.path.join(DATA_DIR, "faiss_index.bin")
META_FILE     = os.path.join(DATA_DIR, "metadata.pkl")
DOCS_LOG      = os.path.join(DATA_DIR, "uploaded_docs.json")
LOG_FILE      = os.path.join(DATA_DIR, "chat_logs.json")
SETTINGS_FILE = os.path.join(DATA_DIR, "settings.json")

CHUNK_SIZE    = 800
TOP_K         = 5
DEFAULT_MODEL = "llama3"
EMBED_MODEL   = "all-MiniLM-L6-v2"
MAX_TOKENS    = 2000
BASE_URL      = "https://thehelpers.vercel.app"
DOWNLOAD_DIR  = os.path.join(
    os.path.expanduser("~"), "Downloads", "TheHelper"
)

# Populated at runtime from Ollama — never a hardcoded list
INSTALLED_MODELS: list = []
FALLBACK_MODELS:  list = ["llama3"]


# ============================================================
#  THEME
# ============================================================

class Theme:
    BG        = "#0f0f0f"
    BG2       = "#1a1a1a"
    BG3       = "#222222"
    BG4       = "#2a2a2a"
    ACCENT    = "#e8c547"
    ACCENT2   = "#4fc3f7"
    FG        = "#e0e0e0"
    FG_DIM    = "#777777"
    RED       = "#ef5350"
    GREEN     = "#66bb6a"
    ORANGE    = "#ffa726"
    BORDER    = "#2e2e2e"
    FONT_MONO = ("Consolas", 10)
    FONT_UI   = ("Segoe UI", 10)
    FONT_SM   = ("Segoe UI", 8)
    FONT_H1   = ("Segoe UI", 16, "bold")
    FONT_H2   = ("Segoe UI", 12, "bold")
    FONT_CODE = ("Consolas", 9)


T = Theme()


# ============================================================
#  SETTINGS
# ============================================================

def save_settings(data: dict):
    try:
        json.dump(
            data,
            open(SETTINGS_FILE, "w", encoding="utf-8"),
            indent=2,
        )
    except Exception:
        pass


def load_settings() -> dict:
    if os.path.exists(SETTINGS_FILE):
        try:
            return json.load(open(SETTINGS_FILE, "r", encoding="utf-8"))
        except Exception:
            pass
    return {}


# ============================================================
#  DOCUMENT PARSING
# ============================================================

def extract_pdf(path: str) -> str:
    if not HAS_PDF:
        return ""
    text = ""
    try:
        reader = PyPDF2.PdfReader(path)
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    except Exception as e:
        text = f"[PDF error: {e}]"
    return text


def extract_pptx(path: str) -> str:
    if not HAS_PPTX:
        return ""
    text = ""
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {i} ---\n"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            text += line + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [
                            c.text.strip()
                            for c in row.cells
                            if c.text.strip()
                        ]
                        if cells:
                            text += " | ".join(cells) + "\n"
            if slide.has_notes_slide:
                notes = (
                    slide.notes_slide.notes_text_frame.text.strip()
                )
                if notes:
                    text += f"Notes: {notes}\n"
    except Exception as e:
        text = f"[PPTX error: {e}]"
    return text


def extract_docx(path: str) -> str:
    if not HAS_DOCX:
        return ""
    text = ""
    try:
        doc = docx.Document(path)
        for para in doc.paragraphs:
            line = para.text.strip()
            if line:
                text += line + "\n"
        for table in doc.tables:
            for row in table.rows:
                cells = [
                    c.text.strip()
                    for c in row.cells
                    if c.text.strip()
                ]
                if cells:
                    text += " | ".join(cells) + "\n"
    except Exception as e:
        text = f"[DOCX error: {e}]"
    return text


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext in (".docx", ".doc"):
        return extract_docx(path)
    if ext in (".md", ".txt"):
        try:
            return open(path, encoding="utf-8", errors="ignore").read()
        except Exception:
            return ""
    return ""


def get_supported_filetypes():
    types = [("Markdown", "*.md"), ("Text", "*.txt")]
    if HAS_PDF:
        types.append(("PDF", "*.pdf"))
    if HAS_PPTX:
        types.append(("PowerPoint", "*.pptx"))
    if HAS_DOCX:
        types.append(("Word", "*.docx *.doc"))
    all_ext = " ".join(t[1] for t in types)
    return [("All Supported", all_ext)] + types


def dep_status_text() -> str:
    return "\n".join([
        "  PyPDF2      : " + (
            "installed" if HAS_PDF
            else "MISSING  pip install PyPDF2"
        ),
        "  python-pptx : " + (
            "installed" if HAS_PPTX
            else "MISSING  pip install python-pptx"
        ),
        "  python-docx : " + (
            "installed" if HAS_DOCX
            else "MISSING  pip install python-docx"
        ),
        "  FAISS+Ollama: " + (
            "installed" if VAULT_AVAILABLE
            else "MISSING  pip install faiss-cpu sentence-transformers ollama"
        ),
        "  Qt Browser  : " + (
            "PySide6/PyQt6" if HAS_QT_BROWSER
            else "MISSING  pip install PySide6"
        ),
    ])


# ============================================================
#  DOCS LOG
# ============================================================

def load_docs_log() -> list:
    if os.path.exists(DOCS_LOG):
        try:
            return json.load(open(DOCS_LOG, "r", encoding="utf-8"))
        except Exception:
            pass
    return []


def save_docs_log(docs: list):
    json.dump(docs, open(DOCS_LOG, "w", encoding="utf-8"), indent=2)


# ============================================================
#  OLLAMA MODEL DISCOVERY
#  Handles both old (dict) and new (ListResponse object) client APIs
# ============================================================

def fetch_installed_models() -> tuple:
    """
    Query the local Ollama daemon for pulled models.

    Returns
    -------
    (names, error)
        names : list[str]  — model name strings, e.g. ["llama3:latest"]
        error : str | None — human-readable error, or None on success
    """
    try:
        import ollama as _ol

        result = _ol.list()

        # ── Normalise response across client versions ─────────────────
        # New client (>=0.2)  → ListResponse object  → .models attribute
        # Old client          → plain dict            → ["models"] key
        if hasattr(result, "models"):
            model_list = result.models
        elif isinstance(result, dict):
            model_list = result.get("models", [])
        else:
            model_list = list(result) if result else []

        names = []
        for m in model_list:
            # New client  → Model object with .model  (full tag)
            #            or .name attribute
            # Old client  → dict with "name" key
            if hasattr(m, "model") and m.model:
                name = m.model.strip()
            elif hasattr(m, "name") and m.name:
                name = m.name.strip()
            elif isinstance(m, dict):
                name = m.get("name", "").strip()
            else:
                name = str(m).strip()

            if name and name not in names:
                names.append(name)

        if names:
            return names, None

        return [], (
            "No models installed yet.\n"
            "Run:  ollama pull llama3"
        )

    except Exception as e:
        return [], f"Cannot reach Ollama: {e}\nRun:  ollama serve"


# ============================================================
#  VAULT CORE
# ============================================================

_embed_model = None
_index       = None
_metadata    = []
_model       = DEFAULT_MODEL


def get_model() -> str:
    return _model


def set_model(name: str):
    global _model
    _model = name


def vault_loaded() -> bool:
    return _index is not None and _embed_model is not None


def _chunk_text(
    text: str, filename: str, filepath: str, source: str
) -> list:
    chunks = []
    for i in range(0, len(text), CHUNK_SIZE):
        chunk = text[i : i + CHUNK_SIZE].strip()
        if chunk:
            chunks.append((
                chunk,
                {
                    "text":      chunk,
                    "path":      filepath,
                    "file":      filename,
                    "source":    source,
                    "chunk_idx": i // CHUNK_SIZE,
                },
            ))
    return chunks


def load_vault(log_fn=None) -> bool:
    global _embed_model, _index, _metadata
    if not VAULT_AVAILABLE:
        if log_fn:
            log_fn("Vault dependencies not installed.")
        return False
    try:
        if log_fn:
            log_fn("Loading embedding model...")
        _embed_model = SentenceTransformer(EMBED_MODEL)
        if os.path.exists(INDEX_FILE) and os.path.exists(META_FILE):
            if log_fn:
                log_fn("Loading existing FAISS index...")
            _index    = faiss.read_index(INDEX_FILE)
            _metadata = pickle.load(open(META_FILE, "rb"))
            if log_fn:
                log_fn(f"Loaded {_index.ntotal} vectors.")
        else:
            rebuild_index(log_fn)
        return True
    except Exception as e:
        if log_fn:
            log_fn(f"Vault load error: {e}")
        return False


def rebuild_index(log_fn=None) -> bool:
    global _index, _metadata
    if _embed_model is None:
        if log_fn:
            log_fn("Embed model not loaded.")
        return False

    if log_fn:
        log_fn("Scanning documents...")

    docs, meta = [], []

    if os.path.exists(NOTES_FOLDER):
        for root_dir, _, files in os.walk(NOTES_FOLDER):
            for f in files:
                if f.endswith(".md"):
                    path = os.path.join(root_dir, f)
                    try:
                        text = open(
                            path, encoding="utf-8", errors="ignore"
                        ).read()
                        chunks = _chunk_text(text, f, path, "obsidian")
                        docs.extend(c[0] for c in chunks)
                        meta.extend(c[1] for c in chunks)
                    except Exception:
                        pass

    if os.path.exists(UPLOADS_DIR):
        for f in os.listdir(UPLOADS_DIR):
            path = os.path.join(UPLOADS_DIR, f)
            if not os.path.isfile(path):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext not in (
                ".pdf", ".pptx", ".docx", ".doc", ".md", ".txt"
            ):
                continue
            try:
                text = extract_text(path)
                if text.strip():
                    chunks = _chunk_text(text, f, path, "uploaded")
                    docs.extend(c[0] for c in chunks)
                    meta.extend(c[1] for c in chunks)
                    if log_fn:
                        log_fn(f"  + {f} ({len(chunks)} chunks)")
            except Exception as e:
                if log_fn:
                    log_fn(f"  ! {f}: {e}")

    if not docs:
        if log_fn:
            log_fn("No content found. Upload some documents first.")
        _index    = None
        _metadata = []
        return False

    if log_fn:
        log_fn(f"Encoding {len(docs)} chunks...")
    emb    = _embed_model.encode(docs).astype("float32")
    _index = faiss.IndexFlatL2(emb.shape[1])
    _index.add(emb)
    _metadata = meta
    faiss.write_index(_index, INDEX_FILE)
    pickle.dump(_metadata, open(META_FILE, "wb"))
    if log_fn:
        log_fn(f"Index built: {len(docs)} chunks.")
    return True


def add_to_index(filepath: str, log_fn=None) -> bool:
    global _index, _metadata
    if _embed_model is None:
        if log_fn:
            log_fn("Vault not loaded.")
        return False
    fname = os.path.basename(filepath)
    text  = extract_text(filepath)
    if not text.strip():
        if log_fn:
            log_fn(f"No text extracted from {fname}")
        return False
    chunks   = _chunk_text(text, fname, filepath, "uploaded")
    new_docs = [c[0] for c in chunks]
    new_meta = [c[1] for c in chunks]
    emb      = _embed_model.encode(new_docs).astype("float32")
    if _index is None:
        _index    = faiss.IndexFlatL2(emb.shape[1])
        _metadata = []
    _index.add(emb)
    _metadata.extend(new_meta)
    faiss.write_index(_index, INDEX_FILE)
    pickle.dump(_metadata, open(META_FILE, "wb"))
    if log_fn:
        log_fn(f"Indexed {len(chunks)} chunks from {fname}")
    return True


def retrieve(query: str) -> list:
    if _index is None or _embed_model is None:
        return []
    q      = _embed_model.encode([query]).astype("float32")
    _, ids = _index.search(q, TOP_K)
    results = []
    for i in ids[0]:
        if 0 <= i < len(_metadata):
            m = _metadata[i]
            results.append({
                "text":   m["text"],
                "file":   m.get("file", "?"),
                "source": m.get("source", "?"),
            })
    return results


def build_messages(context: str, question: str) -> list:
    return [
        {
            "role": "system",
            "content": (
                "You are a knowledgeable assistant. "
                "Answer fully and clearly based on the context provided.\n\n"
                "Structure your response with:\n"
                "1. Clear explanation\n"
                "2. Step-by-step breakdown (if applicable)\n"
                "3. Code examples (if applicable)\n\n"
                "Context:\n" + context
            ),
        },
        {"role": "user", "content": question},
    ]


def stream_response(messages: list, on_token) -> str:
    full   = ""
    stream = ollama.chat(
        model=_model,
        messages=messages,
        stream=True,
        options={"num_predict": MAX_TOKENS},
    )
    for chunk in stream:
        token = chunk["message"]["content"]
        on_token(token)
        full += token
    return full


def save_chat_log(query: str, answer: str):
    logs = []
    if os.path.exists(LOG_FILE):
        try:
            logs = json.load(open(LOG_FILE, "r", encoding="utf-8"))
        except Exception:
            pass
    logs.append({
        "timestamp": datetime.now().isoformat(),
        "model":     _model,
        "query":     query,
        "answer":    answer,
    })
    json.dump(logs, open(LOG_FILE, "w", encoding="utf-8"), indent=2)


# ============================================================
#  BROWSER SCRIPT  (spawned as subprocess)
# ============================================================

_DL   = DOWNLOAD_DIR.replace("\\", "\\\\")
_HOME = BASE_URL

BROWSER_SCRIPT = f'''
import sys, os
start = sys.argv[1] if len(sys.argv) > 1 else "{_HOME}"

try:
    from PySide6.QtWidgets import (QApplication, QMainWindow,
                                   QToolBar, QLineEdit, QProgressBar)
    from PySide6.QtWebEngineWidgets import QWebEngineView
    from PySide6.QtCore import QUrl
    from PySide6.QtGui import QAction
except ImportError:
    try:
        from PyQt6.QtWidgets import (QApplication, QMainWindow,
                                     QToolBar, QLineEdit, QProgressBar)
        from PyQt6.QtWebEngineWidgets import QWebEngineView
        from PyQt6.QtCore import QUrl
        from PyQt6.QtGui import QAction
    except ImportError:
        sys.exit(1)

class Browser(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("TheHelper Browser")
        self.resize(1200, 800)
        self.view = QWebEngineView()
        bar = QToolBar()
        bar.setMovable(False)
        self.addToolBar(bar)
        for label, fn in [("<", self.view.back),
                          (">", self.view.forward),
                          ("R", self.view.reload)]:
            a = QAction(label, self)
            a.triggered.connect(fn)
            bar.addAction(a)
        self.addr = QLineEdit()
        self.addr.returnPressed.connect(self.go)
        bar.addWidget(self.addr)
        go_a = QAction("Go", self)
        go_a.triggered.connect(self.go)
        bar.addAction(go_a)
        self.prog = QProgressBar()
        self.prog.setMaximumWidth(120)
        self.prog.setMaximumHeight(16)
        bar.addWidget(self.prog)
        self.view.urlChanged.connect(
            lambda u: self.addr.setText(u.toString()))
        self.view.loadProgress.connect(self.prog.setValue)
        self.view.loadFinished.connect(lambda _: self.prog.setValue(0))
        self.setCentralWidget(self.view)
        os.makedirs("{_DL}", exist_ok=True)
        self.view.setUrl(QUrl(start))

    def go(self):
        u = self.addr.text().strip()
        if not u.startswith("http"):
            u = "https://" + u
        self.view.setUrl(QUrl(u))

app = QApplication(sys.argv)
win = Browser()
win.show()
sys.exit(app.exec())
'''


# ============================================================
#  WIDGET HELPERS
# ============================================================

def make_frame(parent, bg=None, **kw) -> tk.Frame:
    return tk.Frame(parent, bg=bg or T.BG, **kw)


def make_label(
    parent, text="", font=None, fg=None, bg=None, **kw
) -> tk.Label:
    return tk.Label(
        parent, text=text,
        font=font or T.FONT_UI,
        fg=fg   or T.FG,
        bg=bg   or T.BG,
        **kw,
    )


def make_button(
    parent, text, command, style="normal", **kw
) -> ttk.Button:
    s = "Accent.TButton" if style == "accent" else "TButton"
    return ttk.Button(parent, text=text, command=command, style=s, **kw)


def make_scrolled_text(parent, **kw) -> scrolledtext.ScrolledText:
    defaults = dict(
        wrap="word",
        bg=T.BG2, fg=T.FG,
        insertbackground=T.ACCENT,
        relief="flat", borderwidth=0,
        selectbackground=T.ACCENT,
        selectforeground=T.BG,
        font=T.FONT_MONO,
    )
    defaults.update(kw)
    return scrolledtext.ScrolledText(parent, **defaults)


def hsep(parent, color=None):
    tk.Frame(
        parent, bg=color or T.BORDER, height=1
    ).pack(fill="x")


# ============================================================
#  TTK STYLES
# ============================================================

def apply_styles(root: tk.Tk):
    s = ttk.Style(root)
    s.theme_use("clam")

    s.configure(".",
        background=T.BG, foreground=T.FG,
        fieldbackground=T.BG2, bordercolor=T.BORDER,
        troughcolor=T.BG3, relief="flat", font=T.FONT_UI,
    )
    s.configure("TNotebook",
        background=T.BG, borderwidth=0,
        tabmargins=[0, 0, 0, 0],
    )
    s.configure("TNotebook.Tab",
        background=T.BG2, foreground=T.FG_DIM,
        padding=[18, 7], font=T.FONT_UI, borderwidth=0,
    )
    s.map("TNotebook.Tab",
        background=[("selected", T.BG3), ("active", T.BG3)],
        foreground=[("selected", T.ACCENT), ("active", T.FG)],
    )
    s.configure("TFrame",  background=T.BG)
    s.configure("TLabel",  background=T.BG, foreground=T.FG)
    s.configure("TButton",
        background=T.BG3, foreground=T.FG,
        relief="flat", borderwidth=0, padding=[10, 5],
    )
    s.map("TButton",
        background=[("active", T.BG4), ("pressed", T.BG4)],
        foreground=[("active", T.FG)],
    )
    s.configure("Accent.TButton",
        background=T.ACCENT, foreground=T.BG,
        font=(*T.FONT_UI[:2], "bold"),
        relief="flat", padding=[12, 6],
    )
    s.map("Accent.TButton",
        background=[("active", "#f5d76e"), ("pressed", "#d4a800")],
        foreground=[("active", T.BG)],
    )
    s.configure("Arrow.TButton",
        background=T.BG4, foreground=T.FG,
        relief="flat", borderwidth=0, padding=[4, 5],
    )
    s.map("Arrow.TButton",
        background=[("active", T.ACCENT), ("pressed", T.ACCENT)],
        foreground=[("active", T.BG)],
    )
    s.configure("TEntry",
        fieldbackground=T.BG2, foreground=T.FG,
        insertcolor=T.ACCENT, relief="flat", padding=6,
    )
    s.configure("TCombobox",
        fieldbackground=T.BG2, background=T.BG3,
        foreground=T.FG, arrowcolor=T.ACCENT,
        selectbackground=T.BG3, selectforeground=T.ACCENT,
    )
    s.map("TCombobox",
        fieldbackground=[("readonly", T.BG2)],
        foreground=[("readonly", T.FG)],
        selectbackground=[("readonly", T.BG3)],
        selectforeground=[("readonly", T.ACCENT)],
    )
    s.configure("TScrollbar",
        background=T.BG3, troughcolor=T.BG2,
        arrowcolor=T.FG_DIM, relief="flat", borderwidth=0,
    )
    s.configure("TProgressbar",
        background=T.ACCENT, troughcolor=T.BG2,
        borderwidth=0, relief="flat",
    )

    root.option_add("*Listbox.background",       T.BG2)
    root.option_add("*Listbox.foreground",       T.FG)
    root.option_add("*Listbox.selectBackground", T.ACCENT)
    root.option_add("*Listbox.selectForeground", T.BG)
    root.option_add("*Listbox.relief",           "flat")
    root.option_add("*Listbox.borderWidth",      0)
    root.option_add("*TCombobox*Listbox.background",       T.BG2)
    root.option_add("*TCombobox*Listbox.foreground",       T.FG)
    root.option_add("*TCombobox*Listbox.selectBackground", T.ACCENT)
    root.option_add("*TCombobox*Listbox.selectForeground", T.BG)


# ============================================================
#  TOAST NOTIFICATIONS
# ============================================================

class Toast:
    def __init__(self, root: tk.Tk):
        self.root = root
        self._win = None
        self._job = None

    def show(
        self,
        message: str,
        kind: str = "info",
        duration: int = 3000,
    ):
        color = {
            "info":    T.ACCENT2,
            "success": T.GREEN,
            "warning": T.ORANGE,
            "error":   T.RED,
        }.get(kind, T.FG)

        if self._win:
            try:
                self._win.destroy()
            except Exception:
                pass
        if self._job:
            try:
                self.root.after_cancel(self._job)
            except Exception:
                pass

        win = tk.Toplevel(self.root)
        win.overrideredirect(True)
        win.attributes("-topmost", True)
        win.configure(bg=T.BG3)
        tk.Label(
            win, text=message,
            bg=T.BG3, fg=color,
            font=T.FONT_UI, padx=20, pady=10,
        ).pack()

        self.root.update_idletasks()
        rx = self.root.winfo_x() + self.root.winfo_width()
        ry = self.root.winfo_y() + self.root.winfo_height()
        win.update_idletasks()
        win.geometry(
            f"+{rx - win.winfo_width() - 20}"
            f"+{ry - win.winfo_height() - 40}"
        )
        self._win = win
        self._job = self.root.after(duration, self._dismiss)

    def _dismiss(self):
        if self._win:
            try:
                self._win.destroy()
            except Exception:
                pass
            self._win = None


# ============================================================
#  CHAT TAB
# ============================================================

class ChatTab:
    def __init__(self, notebook: ttk.Notebook, app: "VaultApp"):
        self.app = app
        self.tab = ttk.Frame(notebook)
        notebook.add(self.tab, text="  Chat  ")
        self._build()

    # ── layout ───────────────────────────────────────────────────────

    def _build(self):
        # Output area
        out_wrap = make_frame(self.tab)
        out_wrap.pack(fill="both", expand=True, padx=16, pady=(12, 0))

        self.out = make_scrolled_text(out_wrap, state="disabled")
        self.out.pack(fill="both", expand=True)

        for tag, cfg in {
            "accent":  dict(foreground=T.ACCENT),
            "accent2": dict(foreground=T.ACCENT2),
            "dim":     dict(foreground=T.FG_DIM),
            "green":   dict(foreground=T.GREEN),
            "red":     dict(foreground=T.RED),
            "orange":  dict(foreground=T.ORANGE),
            "heading": dict(
                foreground=T.ACCENT2,
                font=(*T.FONT_MONO[:1], T.FONT_MONO[1] + 1, "bold"),
            ),
            "source": dict(
                foreground=T.FG_DIM,
                font=("Segoe UI", 8, "italic"),
            ),
            "bold": dict(
                font=(*T.FONT_MONO[:1], T.FONT_MONO[1], "bold")
            ),
            "code": dict(
                foreground=T.ACCENT,
                font=T.FONT_CODE,
                background=T.BG3,
            ),
        }.items():
            self.out.tag_configure(tag, **cfg)

        hsep(out_wrap)

        self._status_var = tk.StringVar(value="")
        tk.Label(
            out_wrap,
            textvariable=self._status_var,
            bg=T.BG, fg=T.FG_DIM,
            font=T.FONT_SM, anchor="w",
        ).pack(fill="x", pady=(2, 0))

        # Input row
        inp_wrap = make_frame(self.tab, pady=8)
        inp_wrap.pack(fill="x", padx=16)

        self.inp = tk.Text(
            inp_wrap, height=4,
            bg=T.BG2, fg=T.FG,
            insertbackground=T.ACCENT,
            font=T.FONT_MONO,
            relief="flat", wrap="word",
            selectbackground=T.ACCENT,
            selectforeground=T.BG,
            padx=8, pady=6,
        )
        self.inp.pack(side="left", fill="both", expand=True)
        self.inp.bind("<Return>",    self._on_enter)
        self.inp.bind("<Control-j>", self._on_ctrl_j)

        # Button column
        btn_col = make_frame(inp_wrap, padx=6)
        btn_col.pack(side="left", fill="y")

        make_button(
            btn_col, "Send", self._send, style="accent"
        ).pack(fill="x", pady=(0, 3))

        # ── Split Upload button ───────────────────────────────────────
        upload_row = make_frame(btn_col)
        upload_row.pack(fill="x", pady=(0, 3))

        self._upload_main = ttk.Button(
            upload_row, text="Upload",
            command=self._upload_dialog,
            style="TButton",
        )
        self._upload_main.pack(side="left", fill="x", expand=True)

        self._upload_arrow = ttk.Button(
            upload_row, text="v",
            width=2,
            command=self._show_vault_menu,
            style="Arrow.TButton",
        )
        self._upload_arrow.pack(side="left")
        # ─────────────────────────────────────────────────────────────

        make_button(btn_col, "Clear", self._clear).pack(fill="x")

        # Hint bar
        hint = make_frame(self.tab)
        hint.pack(fill="x", padx=16, pady=(0, 6))
        make_label(
            hint,
            text=(
                "Enter = Send   Ctrl+J = New line   "
                "Upload [v] = pick from vault"
            ),
            font=T.FONT_SM, fg=T.FG_DIM,
        ).pack(side="left")
        self._info_var = tk.StringVar(value="")
        tk.Label(
            hint,
            textvariable=self._info_var,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(side="right")

    # ── keyboard events ───────────────────────────────────────────────

    def _on_enter(self, _e):
        self._send()
        return "break"

    def _on_ctrl_j(self, _e):
        self.inp.insert("insert", "\n")
        return "break"

    # ── send / query ─────────────────────────────────────────────────

    def _send(self):
        query = self.inp.get("1.0", "end").strip()
        if not query:
            return
        self.inp.delete("1.0", "end")
        threading.Thread(
            target=self._run, args=(query,), daemon=True
        ).start()

    def _run(self, query: str):
        if not vault_loaded():
            self.write(
                "\n[Vault not ready - upload documents first]\n",
                "red",
            )
            return

        topics = [
            ln.strip("-* ")
            for ln in query.splitlines()
            if ln.strip()
        ]
        if len(topics) <= 1:
            topics = [query]

        full_answer = ""
        for topic in topics:
            self.write(f"\n>> {topic}\n", "heading")
            results = retrieve(topic)
            if not results:
                self.write("  No relevant context found.\n", "dim")
                continue
            sources = sorted({r["file"] for r in results})
            self.write(
                "  Sources: " + ", ".join(sources) + "\n", "source"
            )
            context  = "\n\n".join(r["text"] for r in results)
            messages = build_messages(context, topic)
            try:
                answer = stream_response(
                    messages, lambda t: self.write(t)
                )
                full_answer += f"\n\n## {topic}\n{answer}"
            except Exception as e:
                self.write(f"\n[Error: {e}]\n", "red")

        if full_answer:
            save_chat_log(query, full_answer)
            self.write("\n\n--- done ---\n", "dim")

        self.app.root.after(0, self.update_info)

    # ── upload: file dialog ──────────────────────────────────────────

    def _upload_dialog(self):
        paths = filedialog.askopenfilenames(
            title="Upload documents",
            filetypes=get_supported_filetypes(),
        )
        if paths:
            self.app.docs_tab.process_uploads(
                list(paths), source="chat"
            )

    # ── upload: vault-docs dropdown ──────────────────────────────────

    def _show_vault_menu(self):
        """
        Pop a menu below the [v] arrow button listing every file
        already stored in the Documents tab.  Selecting one indexes
        it directly from its stored path — no file dialog.
        """
        docs = load_docs_log()

        menu = tk.Menu(
            self.app.root,
            tearoff=0,
            bg=T.BG2, fg=T.FG,
            activebackground=T.ACCENT,
            activeforeground=T.BG,
            font=T.FONT_UI,
            relief="flat", bd=1,
        )

        _ICONS = {
            ".pdf":  "PDF", ".pptx": "PPT",
            ".docx": "DOC", ".doc":  "DOC",
            ".md":   " MD", ".txt":  "TXT",
        }

        if not docs:
            menu.add_command(
                label="  No documents in vault yet",
                state="disabled",
            )
        else:
            for doc in docs:
                icon  = _ICONS.get(doc.get("ext", ""), "???")
                mark  = " [indexed]" if doc.get("indexed") else ""
                label = f"  [{icon}]  {doc['name']}{mark}"
                path  = doc.get("path", "")
                menu.add_command(
                    label=label,
                    command=lambda p=path: self._index_vault_doc(p),
                )
            menu.add_separator()
            menu.add_command(
                label="  + Add ALL vault docs to context",
                command=self._index_all_vault_docs,
            )

        btn = self._upload_arrow
        x   = btn.winfo_rootx()
        y   = btn.winfo_rooty() + btn.winfo_height()
        menu.tk_popup(x, y)
        menu.grab_release()

    def _index_vault_doc(self, path: str):
        """Index a single already-stored file — no file dialog needed."""
        if not path or not os.path.exists(path):
            self.write(f"\n[File not found: {path}]\n", "red")
            return
        fname = os.path.basename(path)
        self.write(f"\nIndexing from vault: {fname}\n", "accent")

        def _run():
            if not vault_loaded():
                self.write("  Vault not ready.\n", "red")
                return
            ok = add_to_index(
                path, lambda m: self.write(f"  {m}\n", "dim")
            )
            if ok:
                self.write(f"  {fname} ready in context.\n", "green")
                self._mark_indexed(path)
                self.app.root.after(0, self.app.docs_tab.refresh_list)
                self.app.root.after(0, self.update_info)
                self.app.toast.show(f"{fname} indexed.", "success")
            else:
                self.write(f"  Could not index {fname}.\n", "red")
                self.app.toast.show(
                    f"Failed to index {fname}.", "error"
                )

        threading.Thread(target=_run, daemon=True).start()

    def _index_all_vault_docs(self):
        """Index every file in the vault docs log."""
        docs  = load_docs_log()
        paths = [
            d["path"]
            for d in docs
            if os.path.exists(d.get("path", ""))
        ]
        if not paths:
            self.write("\nNo documents found in vault.\n", "red")
            return
        self.write(
            f"\nIndexing {len(paths)} vault document(s)...\n",
            "accent",
        )

        def _run():
            ok_n = fail_n = 0
            for path in paths:
                fname = os.path.basename(path)
                ok = add_to_index(
                    path, lambda m: self.write(f"  {m}\n", "dim")
                )
                if ok:
                    self.write(f"  {fname}\n", "green")
                    self._mark_indexed(path)
                    ok_n += 1
                else:
                    self.write(f"  FAILED: {fname}\n", "red")
                    fail_n += 1
            self.app.root.after(0, self.app.docs_tab.refresh_list)
            self.app.root.after(0, self.update_info)
            summary = f"{ok_n} indexed, {fail_n} failed"
            self.write(
                f"\nDone: {summary}\n",
                "green" if fail_n == 0 else "orange",
            )
            self.app.toast.show(
                f"Vault index: {summary}",
                "success" if fail_n == 0 else "warning",
            )

        threading.Thread(target=_run, daemon=True).start()

    @staticmethod
    def _mark_indexed(path: str):
        docs = load_docs_log()
        for d in docs:
            if d.get("path") == path:
                d["indexed"] = True
        save_docs_log(docs)

    # ── clear output ─────────────────────────────────────────────────

    def _clear(self):
        self.out.configure(state="normal")
        self.out.delete("1.0", "end")
        self.out.configure(state="disabled")

    # ── public helpers ────────────────────────────────────────────────

    def write(self, text: str, tag: str = None):
        def _do():
            self.out.configure(state="normal")
            if tag:
                self.out.insert("end", text, (tag,))
            else:
                self.out.insert("end", text)
            self.out.see("end")
            self.out.configure(state="disabled")
        self.app.root.after(0, _do)

    def set_status(self, msg: str):
        self.app.root.after(0, lambda: self._status_var.set(msg))

    def update_info(self):
        n     = _index.ntotal if _index is not None else 0
        ndocs = len(load_docs_log())
        self._info_var.set(
            f"Chunks: {n}   Docs: {ndocs}   Model: {_model}"
        )


# ============================================================
#  DOCUMENTS TAB
# ============================================================

class DocsTab:
    def __init__(self, notebook: ttk.Notebook, app: "VaultApp"):
        self.app = app
        self.tab = ttk.Frame(notebook)
        notebook.add(self.tab, text="  Documents  ")
        self._build()

    def _build(self):
        # Toolbar
        bar = make_frame(self.tab, bg=T.BG2, pady=7, padx=10)
        bar.pack(fill="x")
        make_button(
            bar, "+ Upload Files",
            self._upload_files, style="accent",
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "+ Upload Folder", self._upload_folder
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "Rebuild Index", self._rebuild_index
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "Open Folder", self._open_folder
        ).pack(side="left")
        make_button(
            bar, "Remove All", self._remove_all
        ).pack(side="right", padx=(4, 0))
        make_button(
            bar, "Remove Selected", self._remove_selected
        ).pack(side="right")

        # Status bar + progress
        self._status_var = tk.StringVar(value="  Ready")
        tk.Label(
            self.tab, textvariable=self._status_var,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM, anchor="w",
        ).pack(fill="x", padx=10, pady=(4, 0))
        self._progress = ttk.Progressbar(self.tab, mode="indeterminate")

        # Dependency status panel
        dep = make_frame(self.tab, bg=T.BG3, padx=14, pady=8)
        dep.pack(fill="x", padx=10, pady=(6, 6))
        make_label(
            dep, text="Library status:",
            font=(*T.FONT_UI[:2], "bold"),
            fg=T.ACCENT, bg=T.BG3,
        ).pack(side="left")
        all_ok = HAS_PDF and HAS_PPTX and HAS_DOCX and VAULT_AVAILABLE
        tk.Label(
            dep, text=dep_status_text(),
            font=T.FONT_CODE, bg=T.BG3,
            fg=T.GREEN if all_ok else T.ORANGE,
            justify="left", anchor="w",
        ).pack(side="left", padx=12)

        # Paned: list | preview
        paned = tk.PanedWindow(
            self.tab, orient="horizontal",
            bg=T.BORDER, sashwidth=4,
            sashrelief="flat", bd=0,
        )
        paned.pack(fill="both", expand=True, padx=10, pady=(0, 10))

        # Left — file list
        left = make_frame(paned)
        paned.add(left, width=360, minsize=200)
        make_label(
            left, text="Uploaded Documents",
            font=T.FONT_H2, fg=T.FG, pady=4, anchor="w",
        ).pack(fill="x", padx=6)
        lbox_wrap = tk.Frame(
            left, bg=T.BG2,
            highlightthickness=1,
            highlightbackground=T.BORDER,
        )
        lbox_wrap.pack(fill="both", expand=True, padx=6)
        self.listbox = tk.Listbox(
            lbox_wrap,
            selectmode="extended",
            activestyle="none",
            font=T.FONT_MONO,
        )
        sb = ttk.Scrollbar(lbox_wrap, command=self.listbox.yview)
        self.listbox.configure(yscrollcommand=sb.set)
        sb.pack(side="right", fill="y")
        self.listbox.pack(fill="both", expand=True)
        self.listbox.bind("<<ListboxSelect>>", self._on_select)
        self.listbox.bind("<Double-Button-1>",  self._on_double_click)
        self._list_info = tk.StringVar(value="")
        tk.Label(
            left, textvariable=self._list_info,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM, anchor="w",
        ).pack(fill="x", padx=6, pady=(3, 0))

        # Right — preview
        right = make_frame(paned)
        paned.add(right, minsize=200)
        make_label(
            right, text="Preview",
            font=T.FONT_H2, fg=T.FG, pady=4, anchor="w",
        ).pack(fill="x", padx=6)
        self.preview = make_scrolled_text(
            right, state="disabled", font=T.FONT_CODE
        )
        self.preview.pack(
            fill="both", expand=True, padx=6, pady=(0, 6)
        )
        for tag, cfg in [
            ("heading", dict(
                foreground=T.ACCENT,
                font=(*T.FONT_MONO[:1], 11, "bold"),
            )),
            ("meta", dict(foreground=T.FG_DIM, font=("Segoe UI", 8))),
            ("ok",   dict(foreground=T.GREEN)),
            ("warn", dict(foreground=T.ORANGE)),
            ("err",  dict(foreground=T.RED)),
        ]:
            self.preview.tag_configure(tag, **cfg)

        self.refresh_list()

    # ── helpers ──────────────────────────────────────────────────────

    def _show_progress(self):
        self._progress.pack(fill="x", padx=10, pady=(0, 4))
        self._progress.start(12)

    def _hide_progress(self):
        self._progress.stop()
        self._progress.pack_forget()

    def _set_status(self, msg: str):
        self.app.root.after(
            0, lambda: self._status_var.set("  " + msg)
        )

    def _update_list_info(self):
        docs = load_docs_log()
        n    = _index.ntotal if _index is not None else 0
        self._list_info.set(
            f"  {len(docs)} file(s)   {n} index chunks"
            "   (* = indexed)"
        )

    # ── toolbar actions ───────────────────────────────────────────────

    def _upload_files(self):
        paths = filedialog.askopenfilenames(
            title="Select files to upload",
            filetypes=get_supported_filetypes(),
        )
        if paths:
            threading.Thread(
                target=self.process_uploads,
                args=(list(paths),),
                daemon=True,
            ).start()

    def _upload_folder(self):
        folder = filedialog.askdirectory(title="Select folder")
        if not folder:
            return
        exts  = {".pdf", ".pptx", ".docx", ".doc", ".md", ".txt"}
        paths = [
            os.path.join(folder, f)
            for f in os.listdir(folder)
            if os.path.splitext(f)[1].lower() in exts
        ]
        if not paths:
            messagebox.showinfo(
                "No files", "No supported files found in that folder."
            )
            return
        threading.Thread(
            target=self.process_uploads, args=(paths,), daemon=True
        ).start()

    def _rebuild_index(self):
        if not VAULT_AVAILABLE:
            messagebox.showerror(
                "Missing", "Vault dependencies not installed."
            )
            return
        self._set_status("Rebuilding index...")
        self.app.root.after(0, self._show_progress)
        threading.Thread(target=self._do_rebuild, daemon=True).start()

    def _do_rebuild(self):
        for f in (INDEX_FILE, META_FILE):
            if os.path.exists(f):
                try:
                    os.remove(f)
                except Exception:
                    pass
        rebuild_index(self._set_status)
        self.app.root.after(0, self._hide_progress)
        self.app.root.after(0, self._update_list_info)
        self.app.root.after(0, self.app.chat_tab.update_info)
        self._set_status("Index rebuilt.")
        self.app.toast.show("Index rebuilt.", "success")

    def _open_folder(self):
        os.makedirs(UPLOADS_DIR, exist_ok=True)
        if sys.platform == "win32":
            os.startfile(UPLOADS_DIR)
        elif sys.platform == "darwin":
            subprocess.run(["open", UPLOADS_DIR])
        else:
            subprocess.run(["xdg-open", UPLOADS_DIR])

    def _remove_selected(self):
        sel   = self.listbox.curselection()
        docs  = load_docs_log()
        valid = [i for i in sel if i < len(docs)]
        if not valid:
            messagebox.showinfo(
                "Nothing selected",
                "Select one or more files from the list first.",
            )
            return
        names   = [docs[i]["name"] for i in valid]
        preview = "\n".join(names[:12]) + (
            "\n..." if len(names) > 12 else ""
        )
        if not messagebox.askyesno(
            f"Remove {len(valid)} file(s)?",
            f"Delete from disk:\n\n{preview}",
        ):
            return
        for i in sorted(valid, reverse=True):
            path = docs[i].get("path", "")
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass
            docs.pop(i)
        save_docs_log(docs)
        self.refresh_list()
        self._clear_preview()
        threading.Thread(target=self._do_rebuild, daemon=True).start()

    def _remove_all(self):
        docs = load_docs_log()
        if not docs:
            messagebox.showinfo("Empty", "No documents to remove.")
            return
        if not messagebox.askyesno(
            "Remove ALL?",
            f"Delete all {len(docs)} uploaded files?\n"
            "This cannot be undone.",
        ):
            return
        for doc in docs:
            path = doc.get("path", "")
            if os.path.exists(path):
                try:
                    os.remove(path)
                except Exception:
                    pass
        save_docs_log([])
        self.refresh_list()
        self._clear_preview()
        threading.Thread(target=self._do_rebuild, daemon=True).start()

    # ── list events ───────────────────────────────────────────────────

    def _on_select(self, _e=None):
        sel  = self.listbox.curselection()
        docs = load_docs_log()
        if not sel or sel[0] >= len(docs):
            return
        doc  = docs[sel[0]]
        path = doc.get("path", "")

        self.preview.configure(state="normal")
        self.preview.delete("1.0", "end")
        self.preview.insert("end", doc["name"] + "\n", "heading")
        self.preview.insert(
            "end",
            f"Type     : {doc.get('ext', '?')}\n"
            f"Size     : {doc.get('chars', 0):,} chars\n"
            f"Uploaded : {doc.get('uploaded', '?')[:19]}\n"
            f"Indexed  : ",
            "meta",
        )
        self.preview.insert(
            "end",
            "Yes\n" if doc.get("indexed") else "Pending\n",
            "ok" if doc.get("indexed") else "warn",
        )
        self.preview.insert("end", f"Path     : {path}\n", "meta")
        self.preview.insert("end", "\n" + "-" * 44 + "\n\n", "meta")

        if os.path.exists(path):
            text = extract_text(path)
            snip = text[:3000]
            if len(text) > 3000:
                snip += (
                    f"\n\n... [{len(text):,} chars total"
                    " - showing first 3000]"
                )
            self.preview.insert("end", snip)
        else:
            self.preview.insert(
                "end", "(File not found on disk)", "err"
            )
        self.preview.configure(state="disabled")

    def _on_double_click(self, _e=None):
        sel  = self.listbox.curselection()
        docs = load_docs_log()
        if not sel or sel[0] >= len(docs):
            return
        path = docs[sel[0]].get("path", "")
        if not os.path.exists(path):
            messagebox.showwarning(
                "Not found", f"File not found:\n{path}"
            )
            return
        if sys.platform == "win32":
            subprocess.run(["explorer", "/select,", path])
        elif sys.platform == "darwin":
            subprocess.run(["open", "-R", path])
        else:
            subprocess.run(["xdg-open", os.path.dirname(path)])

    def _clear_preview(self):
        self.preview.configure(state="normal")
        self.preview.delete("1.0", "end")
        self.preview.configure(state="disabled")

    # ── public ───────────────────────────────────────────────────────

    def refresh_list(self):
        self.listbox.delete(0, "end")
        docs = load_docs_log()
        if not docs:
            self.listbox.insert("end", "  No documents yet.")
            self.listbox.insert(
                "end", "  Click '+ Upload Files' above."
            )
            self._update_list_info()
            return
        icons = {
            ".pdf": "PDF", ".pptx": "PPT",
            ".docx": "DOC", ".doc": "DOC",
            ".md": " MD", ".txt": "TXT",
        }
        for doc in docs:
            ext   = doc.get("ext", "")
            icon  = icons.get(ext, "???")
            chars = doc.get("chars", 0)
            size  = (
                f" {chars // 1000}k" if chars > 10_000
                else (f" {chars}" if chars else "")
            )
            mark = " *" if doc.get("indexed") else ""
            self.listbox.insert(
                "end",
                f"  [{icon}]  {doc['name']}{size}{mark}",
            )
        self._update_list_info()

    def process_uploads(self, paths: list, source: str = "docs"):
        """Copy files to UPLOADS_DIR, index them, update log."""
        os.makedirs(UPLOADS_DIR, exist_ok=True)
        self.app.root.after(0, self._show_progress)

        docs_log = load_docs_log()
        total    = len(paths)
        success = failed = 0

        for i, path in enumerate(paths, 1):
            fname = os.path.basename(path)
            self._set_status(f"[{i}/{total}] {fname}")

            dest = os.path.join(UPLOADS_DIR, fname)
            try:
                shutil.copy2(path, dest)
            except Exception as e:
                self._set_status(f"Copy failed: {e}")
                if source == "chat":
                    self.app.chat_tab.write(
                        f"  Copy failed: {e}\n", "red"
                    )
                failed += 1
                continue

            text  = extract_text(dest)
            chars = len(text.strip())
            if chars == 0:
                self._set_status(f"No text extracted: {fname}")
                failed += 1
                continue

            indexed = False
            if vault_loaded():
                ok = add_to_index(
                    dest, lambda m: self._set_status(m)
                )
                if ok:
                    success += 1
                    indexed  = True
                else:
                    failed += 1
            else:
                success += 1

            existing = [d for d in docs_log if d["name"] != fname]
            existing.append({
                "name":     fname,
                "path":     dest,
                "ext":      os.path.splitext(fname)[1].lower(),
                "chars":    chars,
                "uploaded": datetime.now().isoformat(),
                "indexed":  indexed,
            })
            docs_log = existing
            save_docs_log(docs_log)

        self.app.root.after(0, self._hide_progress)
        self.app.root.after(0, self.refresh_list)
        self.app.root.after(0, self.app.chat_tab.update_info)

        summary = f"{success} uploaded, {failed} failed"
        self._set_status(f"Done: {summary}")
        self.app.toast.show(
            summary, "success" if failed == 0 else "warning"
        )
        if source == "chat":
            self.app.chat_tab.write(
                f"\nUpload done: {summary}\n",
                "green" if failed == 0 else "orange",
            )


# ============================================================
#  BROWSER TAB
# ============================================================

class BrowserTab:
    def __init__(self, notebook: ttk.Notebook, app: "VaultApp"):
        self.app       = app
        self.tab       = ttk.Frame(notebook)
        self._proc     = None
        self._launched = False
        notebook.add(self.tab, text="  Browser  ")
        self._build()

    def _build(self):
        bar = make_frame(self.tab, bg=T.BG2, pady=7, padx=10)
        bar.pack(fill="x")
        self._url_var = tk.StringVar(value=BASE_URL)
        addr = ttk.Entry(
            bar, textvariable=self._url_var, font=T.FONT_UI
        )
        addr.pack(side="left", fill="x", expand=True, padx=(0, 6))
        addr.bind(
            "<Return>",
            lambda e: self._navigate(self._url_var.get()),
        )
        make_button(
            bar, "Go",
            lambda: self._navigate(self._url_var.get()),
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "Launch", self._launch, style="accent"
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "Reload", self._relaunch
        ).pack(side="left", padx=(0, 4))
        make_button(
            bar, "System Browser",
            lambda: webbrowser.open(self._url_var.get()),
        ).pack(side="right")

        self._status_var = tk.StringVar(
            value="  Click Launch to open the browser"
        )
        tk.Label(
            self.tab, textvariable=self._status_var,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM, anchor="w",
        ).pack(fill="x", padx=10, pady=(4, 0))
        hsep(self.tab)

        self._content = make_frame(self.tab)
        self._content.pack(fill="both", expand=True)
        self._draw_ready()

    def _clear_content(self):
        for w in self._content.winfo_children():
            w.destroy()

    def _draw_ready(self):
        self._clear_content()
        f = make_frame(self._content)
        f.place(relx=0.5, rely=0.5, anchor="center")
        make_label(
            f, text="TheHelper Browser",
            font=T.FONT_H1, fg=T.FG,
        ).pack(pady=(0, 6))
        make_label(
            f, text="Chromium-powered via Qt WebEngine",
            fg=T.FG_DIM,
        ).pack()
        make_label(
            f, text=f"Downloads -> {DOWNLOAD_DIR}",
            fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(pady=(2, 20))
        if HAS_QT_BROWSER:
            make_button(
                f, "Launch Browser",
                self._launch, style="accent",
            ).pack(pady=4)
        else:
            make_label(
                f,
                text=(
                    "Qt WebEngine not installed.\n"
                    "pip install PySide6"
                ),
                fg=T.ORANGE, font=T.FONT_CODE,
            ).pack(pady=8)
            make_button(
                f, "Open in System Browser",
                lambda: webbrowser.open(self._url_var.get()),
                style="accent",
            ).pack(pady=4)

    def _draw_running(self):
        self._clear_content()
        f = make_frame(self._content)
        f.place(relx=0.5, rely=0.5, anchor="center")
        tk.Label(
            f, text="Browser is running",
            font=T.FONT_H1, bg=T.BG, fg=T.GREEN,
        ).pack(pady=(0, 8))
        make_label(
            f,
            text="Use Alt+Tab to switch to the browser window.",
            fg=T.FG_DIM,
        ).pack()
        make_label(
            f,
            text=f"URL: {self._url_var.get()}",
            fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(pady=(4, 20))
        make_button(
            f, "Relaunch", self._relaunch, style="accent"
        ).pack(pady=4)
        make_button(f, "Stop Browser", self._stop).pack(pady=4)

    # ── process management ────────────────────────────────────────────

    def _running(self) -> bool:
        return self._proc is not None and self._proc.poll() is None

    def _stop(self):
        if self._proc:
            try:
                self._proc.terminate()
                self._proc.wait(timeout=3)
            except Exception:
                try:
                    self._proc.kill()
                except Exception:
                    pass
            self._proc = None
        self._launched = False
        self._status_var.set("  Stopped")
        self._draw_ready()

    def _launch(self):
        if self._running():
            self._status_var.set("  Already running - use Alt+Tab")
            return
        if not HAS_QT_BROWSER:
            messagebox.showerror(
                "Missing",
                "Install Qt WebEngine:\n\npip install PySide6",
            )
            return
        url = self._url_var.get().strip() or BASE_URL
        if not url.startswith("http"):
            url = "https://" + url
        self._url_var.set(url)
        self._status_var.set(f"  Launching: {url}")
        flags = (
            subprocess.CREATE_NO_WINDOW
            if sys.platform == "win32"
            else 0
        )
        try:
            self._proc = subprocess.Popen(
                [sys.executable, "-c", BROWSER_SCRIPT, url],
                creationflags=flags,
                stderr=subprocess.PIPE,
            )
        except Exception as e:
            self._proc = None
            self._status_var.set(f"  Launch failed: {e}")
            return
        self._launched = True
        self.app.root.after(1500, self._check_start)

    def _check_start(self):
        if self._proc is None:
            return
        if self._proc.poll() is not None:
            err = ""
            try:
                err = (
                    self._proc.stderr.read()
                    .decode("utf-8", "ignore")
                    .strip()[:200]
                )
            except Exception:
                pass
            self._proc     = None
            self._launched = False
            self._status_var.set(f"  Crash: {err or 'unknown'}")
            self._draw_ready()
        else:
            self._draw_running()
            self._status_var.set(
                f"  Running (PID {self._proc.pid})"
            )
            self._poll()

    def _poll(self):
        if self._running():
            self.app.root.after(600, self._poll)
        else:
            self._proc     = None
            self._launched = False
            self._status_var.set("  Browser closed")
            self.app.root.after(0, self._draw_ready)

    def _navigate(self, url: str):
        if not url.startswith("http"):
            url = "https://" + url
        self._url_var.set(url)
        if self._running():
            self._stop()
        self.app.root.after(200, self._launch)

    def _relaunch(self):
        self._stop()
        self.app.root.after(300, self._launch)

    def auto_launch(self):
        if not self._launched and HAS_QT_BROWSER:
            self._launch()


# ============================================================
#  HISTORY TAB
# ============================================================

class LogTab:
    def __init__(self, notebook: ttk.Notebook, app: "VaultApp"):
        self.app = app
        self.tab = ttk.Frame(notebook)
        notebook.add(self.tab, text="  History  ")
        self._build()

    def _build(self):
        bar = make_frame(self.tab, pady=8)
        bar.pack(fill="x", padx=16)
        make_button(bar, "Refresh",       self.refresh
                    ).pack(side="left", padx=(0, 6))
        make_button(bar, "Export JSON",   self._export
                    ).pack(side="left", padx=(0, 6))
        make_button(bar, "Clear History", self._clear
                    ).pack(side="left")
        self._info = tk.StringVar(value="")
        tk.Label(
            bar, textvariable=self._info,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(side="right")

        self.view = make_scrolled_text(self.tab, state="disabled")
        self.view.pack(
            fill="both", expand=True, padx=16, pady=(0, 12)
        )
        for tag, cfg in [
            ("ts",  dict(foreground=T.FG_DIM, font=("Segoe UI", 8))),
            ("q",   dict(
                foreground=T.ACCENT,
                font=(*T.FONT_MONO[:1], T.FONT_MONO[1], "bold"),
            )),
            ("mdl", dict(
                foreground=T.ACCENT2, font=("Segoe UI", 8)
            )),
            ("ans", dict(foreground=T.FG)),
        ]:
            self.view.tag_configure(tag, **cfg)
        self.refresh()

    def refresh(self):
        self.view.configure(state="normal")
        self.view.delete("1.0", "end")
        if not os.path.exists(LOG_FILE):
            self.view.insert("end", "No history yet.")
            self.view.configure(state="disabled")
            return
        try:
            logs = json.load(open(LOG_FILE, "r", encoding="utf-8"))
        except Exception as e:
            self.view.insert("end", f"Error: {e}")
            self.view.configure(state="disabled")
            return
        self._info.set(f"{len(logs)} entries")
        for entry in reversed(logs[-100:]):
            ts  = entry.get("timestamp", "")[:19]
            mdl = entry.get("model", "?")
            q   = entry.get("query", "")[:160]
            ans = entry.get("answer", "")[:500]
            self.view.insert("end", f"[{ts}]  ", "ts")
            self.view.insert("end", f"model: {mdl}\n", "mdl")
            self.view.insert("end", f"Q: {q}\n", "q")
            self.view.insert("end", f"{ans}\n\n", "ans")
        self.view.configure(state="disabled")

    def _export(self):
        if not os.path.exists(LOG_FILE):
            messagebox.showinfo("Empty", "No history to export.")
            return
        dest = filedialog.asksaveasfilename(
            defaultextension=".json",
            filetypes=[("JSON", "*.json")],
            initialfile=(
                "vault_history_"
                + datetime.now().strftime("%Y%m%d_%H%M%S")
                + ".json"
            ),
        )
        if dest:
            shutil.copy2(LOG_FILE, dest)
            self.app.toast.show(
                f"Exported to {os.path.basename(dest)}", "success"
            )

    def _clear(self):
        if not messagebox.askyesno(
            "Clear history", "Delete all chat history?"
        ):
            return
        json.dump([], open(LOG_FILE, "w", encoding="utf-8"))
        self.refresh()
        self.app.toast.show("History cleared.", "info")


# ============================================================
#  SETTINGS TAB
# ============================================================

class SettingsTab:
    def __init__(self, notebook: ttk.Notebook, app: "VaultApp"):
        self.app = app
        self.tab = ttk.Frame(notebook)
        notebook.add(self.tab, text="  Settings  ")
        self._build()

    def _build(self):
        canvas = tk.Canvas(self.tab, bg=T.BG, highlightthickness=0)
        sb     = ttk.Scrollbar(
            self.tab, orient="vertical", command=canvas.yview
        )
        canvas.configure(yscrollcommand=sb.set)
        sb.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        inner  = make_frame(canvas)
        win_id = canvas.create_window(
            (0, 0), window=inner, anchor="nw"
        )

        def _resize(e=None):
            canvas.configure(scrollregion=canvas.bbox("all"))
            canvas.itemconfig(win_id, width=canvas.winfo_width())

        inner.bind("<Configure>", _resize)
        canvas.bind("<Configure>", _resize)

        pad = dict(padx=30, pady=6)

        def section(title: str):
            make_label(
                inner, text=title,
                font=T.FONT_H2, fg=T.ACCENT,
            ).pack(anchor="w", **pad)
            tk.Frame(inner, bg=T.BORDER, height=1).pack(
                fill="x", padx=30, pady=(0, 6)
            )

        # ── Model ────────────────────────────────────────────────────
        section("Model")

        self._model_var = tk.StringVar(value=_model)
        self._model_var.trace_add("write", self._on_model_change)

        mf = make_frame(inner)
        mf.pack(fill="x", **pad)
        make_label(
            mf, text="Active model:",
            fg=T.FG_DIM, width=20, anchor="w",
        ).pack(side="left")
        self._model_combo = ttk.Combobox(
            mf,
            textvariable=self._model_var,
            values=INSTALLED_MODELS or FALLBACK_MODELS,
            width=28,
            state="normal",
        )
        self._model_combo.pack(side="left")

        rf = make_frame(inner)
        rf.pack(fill="x", padx=30, pady=(0, 4))
        make_button(
            rf, "Refresh Installed Models", self._refresh_models
        ).pack(side="left")
        self._model_info = tk.StringVar(value="")
        tk.Label(
            rf, textvariable=self._model_info,
            bg=T.BG, fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(side="left", padx=10)

        tk.Label(
            inner,
            text=(
                "  Only models you have pulled with"
                " 'ollama pull <name>' appear here."
            ),
            bg=T.BG, fg=T.FG_DIM,
            font=T.FONT_SM, anchor="w",
        ).pack(fill="x", padx=30, pady=(0, 12))

        # ── Appearance ───────────────────────────────────────────────
        section("Appearance")
        af = make_frame(inner)
        af.pack(fill="x", **pad)
        make_label(
            af, text="Accent colour:",
            fg=T.FG_DIM, width=20, anchor="w",
        ).pack(side="left")
        self._swatch = tk.Label(
            af, text="       ",
            bg=T.ACCENT, cursor="hand2",
            padx=4, pady=4,
        )
        self._swatch.pack(side="left")
        self._swatch.bind("<Button-1>", self._pick_accent)
        make_label(
            af, text=" click to change",
            fg=T.FG_DIM, font=T.FONT_SM,
        ).pack(side="left", padx=6)

        # ── Paths ────────────────────────────────────────────────────
        section("Paths")
        for label, path in [
            ("Data directory",    DATA_DIR),
            ("Uploads",           UPLOADS_DIR),
            ("FAISS index",       INDEX_FILE),
            ("Chat log",          LOG_FILE),
            ("Obsidian notes",    NOTES_FOLDER),
            ("Downloads",         DOWNLOAD_DIR),
        ]:
            pf = make_frame(inner)
            pf.pack(fill="x", padx=30, pady=2)
            make_label(
                pf, text=label + ":",
                fg=T.FG_DIM, width=20, anchor="w",
            ).pack(side="left")
            tk.Label(
                pf, text=path,
                bg=T.BG, fg=T.FG,
                font=T.FONT_CODE, anchor="w",
            ).pack(side="left")

        # ── Vault controls ───────────────────────────────────────────
        section("Vault")
        vf = make_frame(inner)
        vf.pack(fill="x", **pad)
        make_button(
            vf, "Reload Vault", self._reload_vault
        ).pack(side="left", padx=(0, 6))
        make_button(
            vf, "Clear Index", self._clear_index
        ).pack(side="left", padx=(0, 6))
        make_button(
            vf, "Open Data Folder", self._open_data
        ).pack(side="left")

        # ── About ────────────────────────────────────────────────────
        section("About")
        for line in [
            "VaultAI - Local RAG chat over your documents",
            f"Python {sys.version.split()[0]}   "
            f"Tkinter {tk.TkVersion}",
            f"FAISS + Ollama : "
            f"{'available' if VAULT_AVAILABLE else 'not installed'}",
            f"Qt Browser     : "
            f"{'available' if HAS_QT_BROWSER else 'not installed'}",
        ]:
            make_label(
                inner, text=line,
                fg=T.FG_DIM, font=T.FONT_SM,
            ).pack(anchor="w", padx=30, pady=1)

        tk.Frame(inner, height=30, bg=T.BG).pack()

    # ── handlers ─────────────────────────────────────────────────────

    def _on_model_change(self, *_):
        name = self._model_var.get().strip()
        if name and name != "Loading...":
            set_model(name)
            self.app.header.update_model(name)

    def _refresh_models(self):
        self._model_info.set("Querying Ollama...")
        threading.Thread(
            target=self._do_refresh_models, daemon=True
        ).start()

    def _do_refresh_models(self):
        names, err = fetch_installed_models()
        if names:
            self.app.root.after(
                0, lambda ns=names: self._apply_model_list(ns)
            )
        else:
            self.app.root.after(
                0,
                lambda e=err: self._model_info.set(
                    e or "No models found."
                ),
            )

    def _apply_model_list(self, names: list):
        global INSTALLED_MODELS
        INSTALLED_MODELS            = names
        self._model_combo["values"] = names
        self._model_info.set(f"{len(names)} model(s) installed")
        self.app.header.update_model_list(names)
        current = get_model()
        if current not in names and names:
            self._model_var.set(names[0])

    def _pick_accent(self, _e=None):
        result = colorchooser.askcolor(
            color=T.ACCENT,
            title="Pick accent colour",
            parent=self.app.root,
        )
        if result and result[1]:
            self._apply_accent(result[1])

    def _apply_accent(self, colour: str):
        T.ACCENT = colour
        self._swatch.configure(bg=colour)
        self.app.header.apply_accent(colour)

        s = ttk.Style()
        s.configure("Accent.TButton",
                    background=colour, foreground=T.BG)
        s.map("Accent.TButton",
              background=[("active", colour), ("pressed", colour)],
              foreground=[("active", T.BG)])
        s.map("TNotebook.Tab",
              foreground=[("selected", colour)])
        s.configure("TCombobox", arrowcolor=colour)
        s.configure("TProgressbar", background=colour)

        for attr, tags in [
            ("chat_tab.out",     [("accent", colour)]),
            ("log_tab.view",     [("q",      colour)]),
            ("docs_tab.preview", [("heading", colour)]),
        ]:
            parts = attr.split(".")
            obj   = self.app
            try:
                for p in parts:
                    obj = getattr(obj, p)
                for tag, col in tags:
                    obj.tag_configure(tag, foreground=col)
            except Exception:
                pass

        save_settings({"accent": colour, "model": _model})
        self.app.toast.show("Accent colour updated.", "info")

    def _reload_vault(self):
        self.app.toast.show("Reloading vault...", "info")
        threading.Thread(
            target=self._do_reload, daemon=True
        ).start()

    def _do_reload(self):
        ok = load_vault(
            lambda m: self.app.chat_tab.write(m + "\n", "dim")
        )
        if ok:
            self.app.chat_tab.write("Vault reloaded.\n", "green")
            self.app.root.after(0, self.app.chat_tab.update_info)
            self.app.root.after(0, self.app.docs_tab.refresh_list)
            self.app.toast.show("Vault reloaded.", "success")
        else:
            self.app.toast.show("Vault reload failed.", "error")

    def _clear_index(self):
        if not messagebox.askyesno(
            "Clear index",
            "Delete FAISS index files?\n"
            "The docs log will be preserved.",
        ):
            return
        global _index, _metadata
        for f in (INDEX_FILE, META_FILE):
            if os.path.exists(f):
                try:
                    os.remove(f)
                except Exception:
                    pass
        _index    = None
        _metadata = []
        self.app.chat_tab.update_info()
        self.app.docs_tab.refresh_list()
        self.app.toast.show("Index cleared.", "info")

    def _open_data(self):
        if sys.platform == "win32":
            os.startfile(DATA_DIR)
        elif sys.platform == "darwin":
            subprocess.run(["open", DATA_DIR])
        else:
            subprocess.run(["xdg-open", DATA_DIR])


# ============================================================
#  HEADER BAR
# ============================================================

class HeaderBar:
    def __init__(self, root: tk.Tk, app: "VaultApp"):
        self.app   = app
        self.frame = make_frame(root, pady=8)
        self.frame.pack(fill="x", padx=14)
        self._build()

    def _build(self):
        # Logo
        self._logo = tk.Label(
            self.frame, text="VAULT AI",
            font=("Courier New", 20, "bold"),
            bg=T.BG, fg=T.ACCENT,
        )
        self._logo.pack(side="left")
        tk.Label(
            self.frame, text="  +  TheHelper",
            font=("Segoe UI", 11),
            bg=T.BG, fg=T.FG_DIM,
        ).pack(side="left", padx=(2, 20))

        # Status label (rightmost)
        self._status_var = tk.StringVar(value="Initialising...")
        tk.Label(
            self.frame,
            textvariable=self._status_var,
            font=T.FONT_SM, bg=T.BG, fg=T.FG_DIM,
        ).pack(side="right", padx=(8, 0))

        tk.Frame(
            self.frame, bg=T.BORDER, width=1
        ).pack(side="right", fill="y", padx=10, pady=4)

        # Model dropdown — starts empty / "Loading..."
        self._model_var = tk.StringVar(value="Loading...")
        self._model_var.trace_add("write", self._on_model)
        self._combo = ttk.Combobox(
            self.frame,
            textvariable=self._model_var,
            values=["Loading..."],
            width=24,
            state="normal",
            font=("Segoe UI", 9),
        )
        self._combo.pack(side="right")
        tk.Label(
            self.frame, text="Model:",
            font=("Segoe UI", 9),
            bg=T.BG, fg=T.FG_DIM,
        ).pack(side="right", padx=(0, 4))

    def _on_model(self, *_):
        name = self._model_var.get().strip()
        if name and name != "Loading...":
            set_model(name)

    def set_status(self, msg: str):
        self.app.root.after(0, lambda: self._status_var.set(msg))

    def update_model(self, name: str):
        self._model_var.set(name)

    def update_model_list(self, names: list):
        """Replace dropdown with verified installed-models list."""
        self._combo["values"] = names
        current = self._model_var.get()
        if (current not in names or current == "Loading...") and names:
            self._model_var.set(names[0])

    def apply_accent(self, colour: str):
        self._logo.configure(fg=colour)


# ============================================================
#  MAIN APPLICATION
# ============================================================

class VaultApp:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("VaultAI")
        self.root.geometry("1300x820")
        self.root.minsize(960, 640)
        self.root.configure(bg=T.BG)

        settings = load_settings()
        if "accent" in settings:
            T.ACCENT = settings["accent"]
        if "model" in settings:
            set_model(settings["model"])

        apply_styles(root)
        self.toast = Toast(root)

        self.header = HeaderBar(root, self)
        tk.Frame(root, bg=T.BORDER, height=1).pack(fill="x")

        self.nb = ttk.Notebook(root)
        self.nb.pack(fill="both", expand=True)

        self.chat_tab     = ChatTab(self.nb, self)
        self.docs_tab     = DocsTab(self.nb, self)
        self.browser_tab  = BrowserTab(self.nb, self)
        self.log_tab      = LogTab(self.nb, self)
        self.settings_tab = SettingsTab(self.nb, self)

        self.nb.bind("<<NotebookTabChanged>>", self._on_tab_change)
        root.protocol("WM_DELETE_WINDOW", self._on_close)

        threading.Thread(
            target=self._boot, daemon=True
        ).start()

    # ── boot sequence ─────────────────────────────────────────────────

    def _boot(self):
        # Step 1: load vault (embedding model + FAISS)
        self.header.set_status("Loading embedding model...")
        self.chat_tab.write("Starting VaultAI...\n", "dim")

        ok = load_vault(
            lambda m: self.chat_tab.write(m + "\n", "dim")
        )

        if ok:
            chunks = _index.ntotal if _index else 0
            self.header.set_status(f"Ready  *  {chunks} chunks")
            self.chat_tab.write(
                "\nVault ready.  "
                "Upload documents and ask questions.\n",
                "green",
            )
            self.root.after(0, self.chat_tab.update_info)
            self.root.after(0, self.docs_tab.refresh_list)
        else:
            self.header.set_status("Vault unavailable")
            self.chat_tab.write(
                "\nVault dependencies missing.\n"
                "pip install faiss-cpu sentence-transformers ollama\n",
                "red",
            )

        # Step 2: discover installed Ollama models (separate concern)
        self._discover_models()

    def _discover_models(self):
        """
        Query Ollama for locally pulled models and populate dropdowns.
        Uses fetch_installed_models() which handles both old and new
        ollama Python client APIs.
        """
        global INSTALLED_MODELS

        names, err = fetch_installed_models()

        if names:
            INSTALLED_MODELS = names
            # Ensure active model is one that actually exists
            current = get_model()
            if current not in names:
                set_model(names[0])
            self.root.after(
                0, lambda ns=names: self._apply_model_list(ns)
            )
        else:
            self.root.after(
                0, lambda e=err: self._handle_no_models(e)
            )

    def _apply_model_list(self, names: list):
        """Push the verified model list to every dropdown widget."""
        self.header.update_model_list(names)
        try:
            self.settings_tab._model_combo["values"] = names
            current = get_model()
            if current not in names and names:
                self.settings_tab._model_var.set(names[0])
        except Exception:
            pass
        self.chat_tab.write(
            "\nOllama models installed: " + ", ".join(names) + "\n",
            "dim",
        )

    def _handle_no_models(self, err: str):
        """Show a helpful message when Ollama has no models or is offline."""
        self.header.update_model_list(FALLBACK_MODELS)
        if err and "reach" in err.lower():
            # Ollama daemon not running
            self.chat_tab.write(
                f"\n{err}\n"
                "Start Ollama first:  ollama serve\n",
                "red",
            )
        else:
            # Ollama running but nothing pulled
            self.chat_tab.write(
                f"\n{err or 'No models found.'}\n"
                "Pull a model:  ollama pull llama3\n",
                "orange",
            )

    # ── tab switching ─────────────────────────────────────────────────

    def _on_tab_change(self, _e=None):
        idx = self.nb.index(self.nb.select())
        if idx == 2:
            self.root.after(400, self.browser_tab.auto_launch)
        elif idx == 3:
            self.log_tab.refresh()

    # ── shutdown ──────────────────────────────────────────────────────

    def _on_close(self):
        self.browser_tab._stop()
        save_settings({"accent": T.ACCENT, "model": _model})
        self.root.destroy()


# ============================================================
#  ENTRY POINT
# ============================================================

def main():
    root = tk.Tk()
    root.withdraw()
    try:
        icon = os.path.join(os.path.dirname(__file__), "icon.ico")
        if os.path.exists(icon):
            root.iconbitmap(icon)
    except Exception:
        pass
    VaultApp(root)
    root.deiconify()
    root.mainloop()


if __name__ == "__main__":
    main()
